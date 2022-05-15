#!/usr/bin/env python3

####################################################
# ERB Generator Python Script v4.0
# Designed and Written by Salvador Melendez
# GUI Created using: PyQt5 UI code generator 5.14.2
# WARNING! Any changes made to this file can
#          damage the functionality of the script
####################################################


import os
import glob
import xml.etree.ElementTree as ET
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PyQt5 import QtCore, QtGui, QtWidgets
from PyQt5.QtWidgets import QApplication, QWidget, QLabel, QMessageBox, QFileDialog
from PyQt5.QtGui import QIcon, QPixmap, QIntValidator
from PyQt5.QtCore import QDate, QTime, QDateTime, Qt


#VARIABLES
order_list = []
findings_file = 'findings.txt'
desktop_dir = os.path.expanduser("~/Desktop")
cwd = os.getcwd()
current_erb = cwd + '/' + 'erb/findings.xml'
data_folder = ''
data_source = ''
findings = {}
event = []
folder_list = ''
f_folder = ''
scope_list = []
diagram_marking = ''
diagrams_list = []
diagrams_delete = []
strengths_list = []
weaknesses_list = []
mitigrations_list = []
deliverables = []


class Ui_MainWindow(object):
    global order_list, data_folder, findings, event, folder_list
    def setupUi(self, MainWindow):
        MainWindow.setObjectName("MainWindow")
        MainWindow.setEnabled(True)
        MainWindow.resize(1345, 791)
        MainWindow.setMinimumSize(QtCore.QSize(1345, 791))
        MainWindow.setMaximumSize(QtCore.QSize(1345, 791))
        self.centralwidget = QtWidgets.QWidget(MainWindow)
        self.centralwidget.setObjectName("centralwidget")
        self.stackedWidget0 = QtWidgets.QStackedWidget(self.centralwidget)
        self.stackedWidget0.setGeometry(QtCore.QRect(0, 0, 1331, 751))
        self.stackedWidget0.setObjectName("stackedWidget0")
        self.page1 = QtWidgets.QWidget()
        self.page1.setObjectName("page1")
        self.startDate = QtWidgets.QDateEdit(self.page1)
        self.startDate.setGeometry(QtCore.QRect(183, 561, 111, 31))
        self.startDate.setDate(QtCore.QDate(2020, 1, 1))
        self.startDate.setObjectName("startDate")
        self.leadName_label = QtWidgets.QLabel(self.page1)
        self.leadName_label.setGeometry(QtCore.QRect(135, 52, 81, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.leadName_label.setFont(font)
        self.leadName_label.setObjectName("leadName_label")
        self.leadName = QtWidgets.QLineEdit(self.page1)
        self.leadName.setGeometry(QtCore.QRect(208, 56, 301, 31))
        self.leadName.setText("")
        self.leadName.setObjectName("leadName")
        self.leadTitle = QtWidgets.QLineEdit(self.page1)
        self.leadTitle.setGeometry(QtCore.QRect(208, 96, 301, 31))
        self.leadTitle.setObjectName("leadTitle")
        self.eventName = QtWidgets.QLineEdit(self.page1)
        self.eventName.setGeometry(QtCore.QRect(139, 516, 371, 31))
        self.eventName.setText("")
        self.eventName.setObjectName("eventName")
        self.leadTitle_label = QtWidgets.QLabel(self.page1)
        self.leadTitle_label.setGeometry(QtCore.QRect(87, 92, 131, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.leadTitle_label.setFont(font)
        self.leadTitle_label.setObjectName("leadTitle_label")
        self.leadOrg_label = QtWidgets.QLabel(self.page1)
        self.leadOrg_label.setGeometry(QtCore.QRect(58, 133, 151, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.leadOrg_label.setFont(font)
        self.leadOrg_label.setObjectName("leadOrg_label")
        self.endDate = QtWidgets.QDateEdit(self.page1)
        self.endDate.setGeometry(QtCore.QRect(183, 607, 111, 31))
        self.endDate.setDate(QtCore.QDate(2020, 1, 1))
        self.endDate.setObjectName("endDate")
        self.startDate_label = QtWidgets.QLabel(self.page1)
        self.startDate_label.setGeometry(QtCore.QRect(64, 558, 121, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.startDate_label.setFont(font)
        self.startDate_label.setObjectName("startDate_label")
        self.endDate_label = QtWidgets.QLabel(self.page1)
        self.endDate_label.setGeometry(QtCore.QRect(76, 604, 121, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.endDate_label.setFont(font)
        self.endDate_label.setObjectName("endDate_label")
        self.dradisfricData_label = QtWidgets.QLabel(self.page1)
        self.dradisfricData_label.setGeometry(QtCore.QRect(820, 249, 331, 41))
        font = QtGui.QFont()
        font.setPointSize(26)
        font.setBold(True)
        font.setWeight(75)
        self.dradisfricData_label.setFont(font)
        self.dradisfricData_label.setObjectName("dradisfricData_label")
        self.eventName_label = QtWidgets.QLabel(self.page1)
        self.eventName_label.setGeometry(QtCore.QRect(65, 512, 81, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.eventName_label.setFont(font)
        self.eventName_label.setObjectName("eventName_label")
        self.leadOrg = QtWidgets.QLineEdit(self.page1)
        self.leadOrg.setGeometry(QtCore.QRect(208, 137, 301, 31))
        self.leadOrg.setObjectName("leadOrg")
        self.event_label = QtWidgets.QLabel(self.page1)
        self.event_label.setGeometry(QtCore.QRect(30, 458, 131, 41))
        font = QtGui.QFont()
        font.setPointSize(26)
        font.setBold(True)
        font.setWeight(75)
        self.event_label.setFont(font)
        self.event_label.setObjectName("event_label")
        self.teamLead_label = QtWidgets.QLabel(self.page1)
        self.teamLead_label.setGeometry(QtCore.QRect(26, 4, 201, 41))
        font = QtGui.QFont()
        font.setPointSize(26)
        font.setBold(True)
        font.setWeight(75)
        self.teamLead_label.setFont(font)
        self.teamLead_label.setObjectName("teamLead_label")
        self.availFilesFolders_label = QtWidgets.QLabel(self.page1)
        self.availFilesFolders_label.setGeometry(QtCore.QRect(820, 288, 281, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.availFilesFolders_label.setFont(font)
        self.availFilesFolders_label.setObjectName("availFilesFolders_label")
        self.listFolders = QtWidgets.QListWidget(self.page1)
        self.listFolders.setGeometry(QtCore.QRect(810, 319, 371, 371))
        self.listFolders.setObjectName("listFolders")
        self.nextButton = QtWidgets.QPushButton(self.page1)
        self.nextButton.setEnabled(False)
        self.nextButton.setGeometry(QtCore.QRect(1190, 700, 131, 41))
        font = QtGui.QFont()
        font.setPointSize(24)
        font.setBold(True)
        font.setWeight(75)
        self.nextButton.setFont(font)
        self.nextButton.setObjectName("nextButton")
        self.logo_label = QtWidgets.QLabel(self.page1)
        self.logo_label.setGeometry(QtCore.QRect(537, 17, 781, 201))
        self.logo_label.setText("")
        self.logo_label.setPixmap(QtGui.QPixmap("logos.jpg"))
        self.logo_label.setScaledContents(True)
        self.logo_label.setObjectName("logo_label")
        self.Classificationframe = QtWidgets.QFrame(self.page1)
        self.Classificationframe.setGeometry(QtCore.QRect(570, 490, 191, 181))
        self.Classificationframe.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.Classificationframe.setFrameShadow(QtWidgets.QFrame.Raised)
        self.Classificationframe.setObjectName("Classificationframe")
        self.unclassButton = QtWidgets.QRadioButton(self.Classificationframe)
        self.unclassButton.setGeometry(QtCore.QRect(20, 10, 131, 25))
        self.unclassButton.setChecked(True)
        self.unclassButton.setObjectName("unclassButton")
        self.secretButton = QtWidgets.QRadioButton(self.Classificationframe)
        self.secretButton.setGeometry(QtCore.QRect(20, 86, 81, 25))
        self.secretButton.setObjectName("secretButton")
        self.topsecretButton = QtWidgets.QRadioButton(self.Classificationframe)
        self.topsecretButton.setGeometry(QtCore.QRect(20, 130, 121, 25))
        self.topsecretButton.setObjectName("topsecretButton")
        self.noforncheckBox = QtWidgets.QCheckBox(self.Classificationframe)
        self.noforncheckBox.setEnabled(False)
        self.noforncheckBox.setGeometry(QtCore.QRect(48, 105, 91, 25))
        self.noforncheckBox.setObjectName("noforncheckBox")
        self.scicheckBox = QtWidgets.QCheckBox(self.Classificationframe)
        self.scicheckBox.setEnabled(False)
        self.scicheckBox.setGeometry(QtCore.QRect(48, 149, 51, 25))
        self.scicheckBox.setObjectName("scicheckBox")
        self.frame = QtWidgets.QFrame(self.Classificationframe)
        self.frame.setGeometry(QtCore.QRect(40, 32, 151, 51))
        self.frame.setFrameShape(QtWidgets.QFrame.NoFrame)
        self.frame.setFrameShadow(QtWidgets.QFrame.Plain)
        self.frame.setLineWidth(0)
        self.frame.setObjectName("frame")
        self.fouoButton = QtWidgets.QRadioButton(self.frame)
        self.fouoButton.setGeometry(QtCore.QRect(9, 25, 141, 25))
        self.fouoButton.setObjectName("fouoButton")
        self.cuiButton = QtWidgets.QRadioButton(self.frame)
        self.cuiButton.setGeometry(QtCore.QRect(9, 0, 61, 25))
        self.cuiButton.setChecked(True)
        self.cuiButton.setObjectName("cuiButton")
        self.mode_label = QtWidgets.QLabel(self.page1)
        self.mode_label.setGeometry(QtCore.QRect(400, 557, 111, 51))
        font = QtGui.QFont()
        font.setPointSize(15)
        font.setBold(True)
        font.setWeight(75)
        self.mode_label.setFont(font)
        self.mode_label.setLayoutDirection(QtCore.Qt.LeftToRight)
        self.mode_label.setObjectName("mode_label")
        self.Modeframe = QtWidgets.QFrame(self.page1)
        self.Modeframe.setGeometry(QtCore.QRect(400, 608, 111, 81))
        self.Modeframe.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.Modeframe.setFrameShadow(QtWidgets.QFrame.Raised)
        self.Modeframe.setObjectName("Modeframe")
        self.lightButton = QtWidgets.QRadioButton(self.Modeframe)
        self.lightButton.setGeometry(QtCore.QRect(20, 10, 81, 25))
        self.lightButton.setChecked(True)
        self.lightButton.setObjectName("lightButton")
        self.darkButton = QtWidgets.QRadioButton(self.Modeframe)
        self.darkButton.setGeometry(QtCore.QRect(20, 40, 81, 25))
        self.darkButton.setChecked(False)
        self.darkButton.setObjectName("darkButton")
        self.classification_label = QtWidgets.QLabel(self.page1)
        self.classification_label.setGeometry(QtCore.QRect(574, 459, 151, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.classification_label.setFont(font)
        self.classification_label.setLayoutDirection(QtCore.Qt.LeftToRight)
        self.classification_label.setObjectName("classification_label")
        self.draftcheckBox = QtWidgets.QCheckBox(self.page1)
        self.draftcheckBox.setGeometry(QtCore.QRect(550, 673, 211, 21))
        self.draftcheckBox.setChecked(True)
        self.draftcheckBox.setObjectName("draftcheckBox")
        self.office_comboBox = QtWidgets.QComboBox(self.page1)
        self.office_comboBox.setGeometry(QtCore.QRect(210, 405, 151, 27))
        self.office_comboBox.setEditable(False)
        self.office_comboBox.setCurrentText("")
        self.office_comboBox.setObjectName("office_comboBox")
        sup_list = ['', 'FCDD-DAC-E         (Jai)', 'FCDD-DAC-M         (Isabel)', 'FCDD-DAC-O         (Bert)', 'FCDD-DAC-R         (Justin)']
        self.office_comboBox.addItems(sup_list)
        self.office_label = QtWidgets.QLabel(self.page1)
        self.office_label.setGeometry(QtCore.QRect(44, 400, 161, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.office_label.setFont(font)
        self.office_label.setLayoutDirection(QtCore.Qt.LeftToRight)
        self.office_label.setObjectName("office_label")
        self.Eventframe = QtWidgets.QFrame(self.page1)
        self.Eventframe.setGeometry(QtCore.QRect(146, 648, 201, 41))
        self.Eventframe.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.Eventframe.setFrameShadow(QtWidgets.QFrame.Raised)
        self.Eventframe.setObjectName("Eventframe")
        self.pmrButton = QtWidgets.QRadioButton(self.Eventframe)
        self.pmrButton.setGeometry(QtCore.QRect(130, 10, 61, 25))
        self.pmrButton.setChecked(False)
        self.pmrButton.setObjectName("pmrButton")
        self.cvpaButton = QtWidgets.QRadioButton(self.Eventframe)
        self.cvpaButton.setGeometry(QtCore.QRect(10, 10, 111, 25))
        self.cvpaButton.setChecked(True)
        self.cvpaButton.setObjectName("cvpaButton")
        self.eventType_label = QtWidgets.QLabel(self.page1)
        self.eventType_label.setGeometry(QtCore.QRect(80, 650, 71, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.eventType_label.setFont(font)
        self.eventType_label.setObjectName("eventType_label")
        self.leadEmail_label = QtWidgets.QLabel(self.page1)
        self.leadEmail_label.setGeometry(QtCore.QRect(126, 173, 81, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.leadEmail_label.setFont(font)
        self.leadEmail_label.setObjectName("leadEmail_label")
        self.leadNIPR = QtWidgets.QLineEdit(self.page1)
        self.leadNIPR.setGeometry(QtCore.QRect(208, 210, 301, 31))
        self.leadNIPR.setObjectName("leadNIPR")
        self.leadNIPR_label = QtWidgets.QLabel(self.page1)
        self.leadNIPR_label.setGeometry(QtCore.QRect(150, 207, 61, 31))
        font = QtGui.QFont()
        font.setPointSize(16)
        font.setBold(True)
        font.setWeight(75)
        self.leadNIPR_label.setFont(font)
        self.leadNIPR_label.setObjectName("leadNIPR_label")
        self.leadSIPR = QtWidgets.QLineEdit(self.page1)
        self.leadSIPR.setGeometry(QtCore.QRect(209, 250, 301, 31))
        self.leadSIPR.setObjectName("leadSIPR")
        self.leadSIPR_label = QtWidgets.QLabel(self.page1)
        self.leadSIPR_label.setGeometry(QtCore.QRect(154, 247, 61, 31))
        font = QtGui.QFont()
        font.setPointSize(16)
        font.setBold(True)
        font.setWeight(75)
        self.leadSIPR_label.setFont(font)
        self.leadSIPR_label.setObjectName("leadSIPR_label")
        self.leadPhone_label = QtWidgets.QLabel(self.page1)
        self.leadPhone_label.setGeometry(QtCore.QRect(117, 283, 91, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.leadPhone_label.setFont(font)
        self.leadPhone_label.setObjectName("leadPhone_label")
        self.leadOffice = QtWidgets.QLineEdit(self.page1)
        self.leadOffice.setGeometry(QtCore.QRect(208, 320, 301, 31))
        self.leadOffice.setObjectName("leadOffice")
        self.leadMobile = QtWidgets.QLineEdit(self.page1)
        self.leadMobile.setGeometry(QtCore.QRect(209, 360, 301, 31))
        self.leadMobile.setObjectName("leadMobile")
        self.leadOffice_label = QtWidgets.QLabel(self.page1)
        self.leadOffice_label.setGeometry(QtCore.QRect(140, 317, 71, 31))
        font = QtGui.QFont()
        font.setPointSize(16)
        font.setBold(True)
        font.setWeight(75)
        self.leadOffice_label.setFont(font)
        self.leadOffice_label.setObjectName("leadOffice_label")
        self.leadMobile_label = QtWidgets.QLabel(self.page1)
        self.leadMobile_label.setGeometry(QtCore.QRect(132, 357, 81, 31))
        font = QtGui.QFont()
        font.setPointSize(16)
        font.setBold(True)
        font.setWeight(75)
        self.leadMobile_label.setFont(font)
        self.leadMobile_label.setObjectName("leadMobile_label")
        self.systemButton = QtWidgets.QPushButton(self.page1)
        self.systemButton.setEnabled(True)
        self.systemButton.setGeometry(QtCore.QRect(570, 319, 191, 91))
        font = QtGui.QFont()
        font.setPointSize(24)
        font.setBold(True)
        font.setWeight(75)
        self.systemButton.setFont(font)
        self.systemButton.setObjectName("systemButton")
        self.stackedWidget0.addWidget(self.page1)
        self.page2 = QtWidgets.QWidget()
        self.page2.setObjectName("page2")
        self.quitButton = QtWidgets.QPushButton(self.page2)
        self.quitButton.setGeometry(QtCore.QRect(1220, 640, 91, 41))
        font = QtGui.QFont()
        font.setPointSize(24)
        font.setBold(True)
        font.setWeight(75)
        self.quitButton.setFont(font)
        self.quitButton.setObjectName("quitButton")
        self.findings_label = QtWidgets.QLabel(self.page2)
        self.findings_label.setGeometry(QtCore.QRect(25, 0, 111, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.findings_label.setFont(font)
        self.findings_label.setObjectName("findings_label")
        self.updateDescButton = QtWidgets.QPushButton(self.page2)
        self.updateDescButton.setGeometry(QtCore.QRect(298, 702, 181, 41))
        font = QtGui.QFont()
        font.setPointSize(14)
        font.setBold(True)
        font.setWeight(75)
        self.updateDescButton.setFont(font)
        self.updateDescButton.setObjectName("updateDescButton")
        self.issuesText = QtWidgets.QPlainTextEdit(self.page2)
        self.issuesText.setGeometry(QtCore.QRect(19, 330, 461, 121))
        self.issuesText.setObjectName("issuesText")
        self.screenshotPreview_label = QtWidgets.QLabel(self.page2)
        self.screenshotPreview_label.setGeometry(QtCore.QRect(550, 352, 231, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.screenshotPreview_label.setFont(font)
        self.screenshotPreview_label.setObjectName("screenshotPreview_label")
        self.imagePreview_label = QtWidgets.QLabel(self.page2)
        self.imagePreview_label.setGeometry(QtCore.QRect(540, 383, 631, 391))
        self.imagePreview_label.setScaledContents(True)
        self.imagePreview_label.setAlignment(QtCore.Qt.AlignCenter)
        self.imagePreview_label.setObjectName("imagePreview_label")
        self.findingDeleteButton = QtWidgets.QPushButton(self.page2)
        self.findingDeleteButton.setGeometry(QtCore.QRect(439, 178, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.findingDeleteButton.setFont(font)
        self.findingDeleteButton.setObjectName("findingDeleteButton")
        self.findingUpButton = QtWidgets.QPushButton(self.page2)
        self.findingUpButton.setGeometry(QtCore.QRect(440, 98, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.findingUpButton.setFont(font)
        self.findingUpButton.setCursor(QtGui.QCursor(QtCore.Qt.ArrowCursor))
        self.findingUpButton.setObjectName("findingUpButton")
        self.findingDownButton = QtWidgets.QPushButton(self.page2)
        self.findingDownButton.setGeometry(QtCore.QRect(440, 128, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.findingDownButton.setFont(font)
        self.findingDownButton.setObjectName("findingDownButton")
        self.screenshots_label = QtWidgets.QLabel(self.page2)
        self.screenshots_label.setGeometry(QtCore.QRect(566, 0, 141, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.screenshots_label.setFont(font)
        self.screenshots_label.setObjectName("screenshots_label")
        self.pptxButton = QtWidgets.QPushButton(self.page2)
        self.pptxButton.setGeometry(QtCore.QRect(1220, 30, 91, 41))
        font = QtGui.QFont()
        font.setPointSize(24)
        font.setBold(True)
        font.setWeight(75)
        self.pptxButton.setFont(font)
        self.pptxButton.setObjectName("pptxButton")
        self.issues_label = QtWidgets.QLabel(self.page2)
        self.issues_label.setGeometry(QtCore.QRect(29, 299, 51, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(True)
        font.setWeight(75)
        self.issues_label.setFont(font)
        self.issues_label.setObjectName("issues_label")
        self.listFindings = QtWidgets.QListWidget(self.page2)
        self.listFindings.setGeometry(QtCore.QRect(15, 30, 421, 181))
        self.listFindings.setDragEnabled(False)
        self.listFindings.setDragDropOverwriteMode(False)
        self.listFindings.setDragDropMode(QtWidgets.QAbstractItemView.NoDragDrop)
        self.listFindings.setDefaultDropAction(QtCore.Qt.IgnoreAction)
        self.listFindings.setObjectName("listFindings")
        self.gobackButton = QtWidgets.QPushButton(self.page2)
        self.gobackButton.setGeometry(QtCore.QRect(1200, 700, 131, 41))
        font = QtGui.QFont()
        font.setPointSize(24)
        font.setBold(True)
        font.setWeight(75)
        self.gobackButton.setFont(font)
        self.gobackButton.setObjectName("gobackButton")
        self.screenshotDeleteButton = QtWidgets.QPushButton(self.page2)
        self.screenshotDeleteButton.setGeometry(QtCore.QRect(1129, 240, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.screenshotDeleteButton.setFont(font)
        self.screenshotDeleteButton.setObjectName("screenshotDeleteButton")
        self.screenshotDownButton = QtWidgets.QPushButton(self.page2)
        self.screenshotDownButton.setGeometry(QtCore.QRect(1130, 190, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.screenshotDownButton.setFont(font)
        self.screenshotDownButton.setObjectName("screenshotDownButton")
        self.screenshotUpButton = QtWidgets.QPushButton(self.page2)
        self.screenshotUpButton.setGeometry(QtCore.QRect(1130, 160, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.screenshotUpButton.setFont(font)
        self.screenshotUpButton.setCursor(QtGui.QCursor(QtCore.Qt.ArrowCursor))
        self.screenshotUpButton.setObjectName("screenshotUpButton")
        self.listScreenshots = QtWidgets.QListWidget(self.page2)
        self.listScreenshots.setGeometry(QtCore.QRect(556, 30, 571, 321))
        self.listScreenshots.setDragEnabled(False)
        self.listScreenshots.setDragDropOverwriteMode(False)
        self.listScreenshots.setDragDropMode(QtWidgets.QAbstractItemView.NoDragDrop)
        self.listScreenshots.setDefaultDropAction(QtCore.Qt.IgnoreAction)
        self.listScreenshots.setObjectName("listScreenshots")
        self.findingName_label = QtWidgets.QLabel(self.page2)
        self.findingName_label.setGeometry(QtCore.QRect(33, 228, 101, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(True)
        font.setWeight(75)
        self.findingName_label.setFont(font)
        self.findingName_label.setObjectName("findingName_label")
        self.findingName = QtWidgets.QLineEdit(self.page2)
        self.findingName.setGeometry(QtCore.QRect(139, 229, 341, 31))
        self.findingName.setText("")
        self.findingName.setPlaceholderText("")
        self.findingName.setObjectName("findingName")
        self.findingHosts_label = QtWidgets.QLabel(self.page2)
        self.findingHosts_label.setGeometry(QtCore.QRect(29, 267, 111, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(True)
        font.setWeight(75)
        self.findingHosts_label.setFont(font)
        self.findingHosts_label.setObjectName("findingHosts_label")
        self.findingHosts = QtWidgets.QLineEdit(self.page2)
        self.findingHosts.setGeometry(QtCore.QRect(139, 269, 341, 31))
        self.findingHosts.setText("")
        self.findingHosts.setPlaceholderText("")
        self.findingHosts.setObjectName("findingHosts")
        self.posture_label = QtWidgets.QLabel(self.page2)
        self.posture_label.setGeometry(QtCore.QRect(19, 450, 61, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(True)
        font.setWeight(75)
        self.posture_label.setFont(font)
        self.posture_label.setObjectName("posture_label")
        self.insiderButton = QtWidgets.QRadioButton(self.page2)
        self.insiderButton.setGeometry(QtCore.QRect(36, 477, 91, 25))
        self.insiderButton.setObjectName("insiderButton")
        self.nearsiderButton = QtWidgets.QRadioButton(self.page2)
        self.nearsiderButton.setGeometry(QtCore.QRect(36, 497, 91, 25))
        self.nearsiderButton.setChecked(True)
        self.nearsiderButton.setObjectName("nearsiderButton")
        self.outsiderButton = QtWidgets.QRadioButton(self.page2)
        self.outsiderButton.setGeometry(QtCore.QRect(36, 517, 91, 25))
        self.outsiderButton.setChecked(False)
        self.outsiderButton.setObjectName("outsiderButton")
        self.mitigationText = QtWidgets.QPlainTextEdit(self.page2)
        self.mitigationText.setGeometry(QtCore.QRect(19, 566, 461, 121))
        self.mitigationText.setObjectName("mitigationText")
        self.mitigation_label = QtWidgets.QLabel(self.page2)
        self.mitigation_label.setGeometry(QtCore.QRect(29, 539, 81, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(True)
        font.setWeight(75)
        self.mitigation_label.setFont(font)
        self.mitigation_label.setObjectName("mitigation_label")
        self.mitigationcheckBox = QtWidgets.QCheckBox(self.page2)
        self.mitigationcheckBox.setGeometry(QtCore.QRect(29, 696, 141, 25))
        self.mitigationcheckBox.setObjectName("mitigationcheckBox")
        self.findingAddButton = QtWidgets.QPushButton(self.page2)
        self.findingAddButton.setGeometry(QtCore.QRect(440, 48, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.findingAddButton.setFont(font)
        self.findingAddButton.setObjectName("findingAddButton")
        self.screenshotAddButton = QtWidgets.QPushButton(self.page2)
        self.screenshotAddButton.setGeometry(QtCore.QRect(1130, 100, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.screenshotAddButton.setFont(font)
        self.screenshotAddButton.setObjectName("screenshotAddButton")
        self.imagePreview_label.raise_()
        self.quitButton.raise_()
        self.findings_label.raise_()
        self.updateDescButton.raise_()
        self.issuesText.raise_()
        self.screenshotPreview_label.raise_()
        self.findingDeleteButton.raise_()
        self.findingUpButton.raise_()
        self.findingDownButton.raise_()
        self.screenshots_label.raise_()
        self.pptxButton.raise_()
        self.issues_label.raise_()
        self.listFindings.raise_()
        self.gobackButton.raise_()
        self.screenshotDeleteButton.raise_()
        self.screenshotDownButton.raise_()
        self.screenshotUpButton.raise_()
        self.listScreenshots.raise_()
        self.findingName_label.raise_()
        self.findingName.raise_()
        self.findingHosts_label.raise_()
        self.findingHosts.raise_()
        self.posture_label.raise_()
        self.insiderButton.raise_()
        self.nearsiderButton.raise_()
        self.outsiderButton.raise_()
        self.mitigationText.raise_()
        self.mitigation_label.raise_()
        self.mitigationcheckBox.raise_()
        self.findingAddButton.raise_()
        self.screenshotAddButton.raise_()
        self.stackedWidget0.addWidget(self.page2)
        self.page = QtWidgets.QWidget()
        self.page.setObjectName("page")
        self.strengthsDeleteButton = QtWidgets.QPushButton(self.page)
        self.strengthsDeleteButton.setGeometry(QtCore.QRect(1289, 139, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.strengthsDeleteButton.setFont(font)
        self.strengthsDeleteButton.setObjectName("strengthsDeleteButton")
        self.strengths_label = QtWidgets.QLabel(self.page)
        self.strengths_label.setGeometry(QtCore.QRect(939, 70, 121, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.strengths_label.setFont(font)
        self.strengths_label.setObjectName("strengths_label")
        self.listStrengths = QtWidgets.QListWidget(self.page)
        self.listStrengths.setGeometry(QtCore.QRect(935, 139, 350, 150))
        self.listStrengths.setDragEnabled(False)
        self.listStrengths.setDragDropOverwriteMode(False)
        self.listStrengths.setDragDropMode(QtWidgets.QAbstractItemView.NoDragDrop)
        self.listStrengths.setDefaultDropAction(QtCore.Qt.IgnoreAction)
        self.listStrengths.setObjectName("listStrengths")
        self.strengthsAddButton = QtWidgets.QPushButton(self.page)
        self.strengthsAddButton.setGeometry(QtCore.QRect(1289, 104, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.strengthsAddButton.setFont(font)
        self.strengthsAddButton.setObjectName("strengthsAddButton")
        self.sutPreview_label = QtWidgets.QLabel(self.page)
        self.sutPreview_label.setGeometry(QtCore.QRect(480, 328, 151, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.sutPreview_label.setFont(font)
        self.sutPreview_label.setObjectName("sutPreview_label")
        self.sutPreview_image = QtWidgets.QLabel(self.page)
        self.sutPreview_image.setGeometry(QtCore.QRect(470, 368, 420, 300))
        self.sutPreview_image.setScaledContents(True)
        self.sutPreview_image.setAlignment(QtCore.Qt.AlignCenter)
        self.sutPreview_image.setObjectName("sutPreview_image")
        self.SUT_label = QtWidgets.QLabel(self.page)
        self.SUT_label.setGeometry(QtCore.QRect(20, 20, 451, 41))
        font = QtGui.QFont()
        font.setPointSize(26)
        font.setBold(True)
        font.setWeight(75)
        self.SUT_label.setFont(font)
        self.SUT_label.setObjectName("SUT_label")
        self.listWeaknesses = QtWidgets.QListWidget(self.page)
        self.listWeaknesses.setGeometry(QtCore.QRect(935, 359, 350, 150))
        self.listWeaknesses.setDragEnabled(False)
        self.listWeaknesses.setDragDropOverwriteMode(False)
        self.listWeaknesses.setDragDropMode(QtWidgets.QAbstractItemView.NoDragDrop)
        self.listWeaknesses.setDefaultDropAction(QtCore.Qt.IgnoreAction)
        self.listWeaknesses.setObjectName("listWeaknesses")
        self.weaknesses_label = QtWidgets.QLabel(self.page)
        self.weaknesses_label.setGeometry(QtCore.QRect(945, 290, 151, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.weaknesses_label.setFont(font)
        self.weaknesses_label.setObjectName("weaknesses_label")
        self.weaknessesAddButton = QtWidgets.QPushButton(self.page)
        self.weaknessesAddButton.setGeometry(QtCore.QRect(1289, 323, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.weaknessesAddButton.setFont(font)
        self.weaknessesAddButton.setObjectName("weaknessesAddButton")
        self.weaknessesDeleteButton = QtWidgets.QPushButton(self.page)
        self.weaknessesDeleteButton.setGeometry(QtCore.QRect(1289, 359, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.weaknessesDeleteButton.setFont(font)
        self.weaknessesDeleteButton.setObjectName("weaknessesDeleteButton")
        self.mitigationsDeleteButton = QtWidgets.QPushButton(self.page)
        self.mitigationsDeleteButton.setGeometry(QtCore.QRect(1289, 579, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.mitigationsDeleteButton.setFont(font)
        self.mitigationsDeleteButton.setObjectName("mitigationsDeleteButton")
        self.mitigationsAddButton = QtWidgets.QPushButton(self.page)
        self.mitigationsAddButton.setGeometry(QtCore.QRect(1289, 543, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.mitigationsAddButton.setFont(font)
        self.mitigationsAddButton.setObjectName("mitigationsAddButton")
        self.listMitigations = QtWidgets.QListWidget(self.page)
        self.listMitigations.setGeometry(QtCore.QRect(935, 579, 350, 150))
        self.listMitigations.setDragEnabled(False)
        self.listMitigations.setDragDropOverwriteMode(False)
        self.listMitigations.setDragDropMode(QtWidgets.QAbstractItemView.NoDragDrop)
        self.listMitigations.setDefaultDropAction(QtCore.Qt.IgnoreAction)
        self.listMitigations.setObjectName("listMitigations")
        self.mitigations_label = QtWidgets.QLabel(self.page)
        self.mitigations_label.setGeometry(QtCore.QRect(945, 510, 221, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.mitigations_label.setFont(font)
        self.mitigations_label.setObjectName("mitigations_label")
        self.inputStrength = QtWidgets.QLineEdit(self.page)
        self.inputStrength.setGeometry(QtCore.QRect(1004, 103, 280, 31))
        self.inputStrength.setText("")
        self.inputStrength.setObjectName("inputStrength")
        self.inputStrength_label = QtWidgets.QLabel(self.page)
        self.inputStrength_label.setGeometry(QtCore.QRect(939, 100, 61, 31))
        font = QtGui.QFont()
        font.setPointSize(16)
        font.setBold(True)
        font.setWeight(75)
        self.inputStrength_label.setFont(font)
        self.inputStrength_label.setObjectName("inputStrength_label")
        self.inputWeakness_label = QtWidgets.QLabel(self.page)
        self.inputWeakness_label.setGeometry(QtCore.QRect(939, 320, 61, 31))
        font = QtGui.QFont()
        font.setPointSize(16)
        font.setBold(True)
        font.setWeight(75)
        self.inputWeakness_label.setFont(font)
        self.inputWeakness_label.setObjectName("inputWeakness_label")
        self.inputWeakness = QtWidgets.QLineEdit(self.page)
        self.inputWeakness.setGeometry(QtCore.QRect(1004, 323, 280, 31))
        self.inputWeakness.setText("")
        self.inputWeakness.setObjectName("inputWeakness")
        self.inputMitigation_label = QtWidgets.QLabel(self.page)
        self.inputMitigation_label.setGeometry(QtCore.QRect(939, 540, 61, 31))
        font = QtGui.QFont()
        font.setPointSize(16)
        font.setBold(True)
        font.setWeight(75)
        self.inputMitigation_label.setFont(font)
        self.inputMitigation_label.setObjectName("inputMitigation_label")
        self.inputMitigation = QtWidgets.QLineEdit(self.page)
        self.inputMitigation.setGeometry(QtCore.QRect(1004, 543, 280, 31))
        self.inputMitigation.setText("")
        self.inputMitigation.setObjectName("inputMitigation")
        self.saveButton = QtWidgets.QPushButton(self.page)
        self.saveButton.setGeometry(QtCore.QRect(1210, 30, 91, 41))
        font = QtGui.QFont()
        font.setPointSize(24)
        font.setBold(True)
        font.setWeight(75)
        self.saveButton.setFont(font)
        self.saveButton.setObjectName("saveButton")
        self.diagramSUT_label = QtWidgets.QLabel(self.page)
        self.diagramSUT_label.setGeometry(QtCore.QRect(490, 120, 171, 31))
        font = QtGui.QFont()
        font.setPointSize(16)
        font.setBold(True)
        font.setWeight(75)
        self.diagramSUT_label.setFont(font)
        self.diagramSUT_label.setObjectName("diagramSUT_label")
        self.figureBrowseButton = QtWidgets.QPushButton(self.page)
        self.figureBrowseButton.setGeometry(QtCore.QRect(655, 122, 71, 31))
        font = QtGui.QFont()
        font.setPointSize(10)
        font.setBold(True)
        font.setWeight(75)
        self.figureBrowseButton.setFont(font)
        self.figureBrowseButton.setObjectName("figureBrowseButton")
        self.inputScope_label = QtWidgets.QLabel(self.page)
        self.inputScope_label.setGeometry(QtCore.QRect(30, 151, 61, 31))
        font = QtGui.QFont()
        font.setPointSize(16)
        font.setBold(True)
        font.setWeight(75)
        self.inputScope_label.setFont(font)
        self.inputScope_label.setObjectName("inputScope_label")
        self.listScope = QtWidgets.QListWidget(self.page)
        self.listScope.setGeometry(QtCore.QRect(26, 240, 350, 150))
        self.listScope.setDragEnabled(False)
        self.listScope.setDragDropOverwriteMode(False)
        self.listScope.setDragDropMode(QtWidgets.QAbstractItemView.NoDragDrop)
        self.listScope.setDefaultDropAction(QtCore.Qt.IgnoreAction)
        self.listScope.setObjectName("listScope")
        self.scope_label = QtWidgets.QLabel(self.page)
        self.scope_label.setGeometry(QtCore.QRect(30, 121, 81, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.scope_label.setFont(font)
        self.scope_label.setObjectName("scope_label")
        self.scopeAddButton = QtWidgets.QPushButton(self.page)
        self.scopeAddButton.setGeometry(QtCore.QRect(380, 155, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.scopeAddButton.setFont(font)
        self.scopeAddButton.setObjectName("scopeAddButton")
        self.inputScope = QtWidgets.QLineEdit(self.page)
        self.inputScope.setGeometry(QtCore.QRect(95, 154, 280, 31))
        self.inputScope.setText("")
        self.inputScope.setObjectName("inputScope")
        self.scopeDeleteButton = QtWidgets.QPushButton(self.page)
        self.scopeDeleteButton.setGeometry(QtCore.QRect(380, 240, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.scopeDeleteButton.setFont(font)
        self.scopeDeleteButton.setObjectName("scopeDeleteButton")
        self.Scopeframe = QtWidgets.QFrame(self.page)
        self.Scopeframe.setGeometry(QtCore.QRect(96, 186, 251, 41))
        self.Scopeframe.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.Scopeframe.setFrameShadow(QtWidgets.QFrame.Raised)
        self.Scopeframe.setObjectName("Scopeframe")
        self.subbulletButton = QtWidgets.QRadioButton(self.Scopeframe)
        self.subbulletButton.setGeometry(QtCore.QRect(90, 10, 151, 25))
        self.subbulletButton.setChecked(False)
        self.subbulletButton.setObjectName("subbulletButton")
        self.bulletButton = QtWidgets.QRadioButton(self.Scopeframe)
        self.bulletButton.setGeometry(QtCore.QRect(10, 10, 71, 25))
        self.bulletButton.setChecked(True)
        self.bulletButton.setObjectName("bulletButton")
        self.Deliverablesframe = QtWidgets.QFrame(self.page)
        self.Deliverablesframe.setGeometry(QtCore.QRect(88, 615, 181, 101))
        self.Deliverablesframe.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.Deliverablesframe.setFrameShadow(QtWidgets.QFrame.Raised)
        self.Deliverablesframe.setObjectName("Deliverablesframe")
        self.reportButton = QtWidgets.QRadioButton(self.Deliverablesframe)
        self.reportButton.setGeometry(QtCore.QRect(10, 33, 71, 25))
        self.reportButton.setChecked(True)
        self.reportButton.setObjectName("reportButton")
        self.memoButton = QtWidgets.QRadioButton(self.Deliverablesframe)
        self.memoButton.setGeometry(QtCore.QRect(10, 10, 121, 25))
        self.memoButton.setChecked(False)
        self.memoButton.setObjectName("memoButton")
        self.inputDays = QtWidgets.QLineEdit(self.Deliverablesframe)
        self.inputDays.setGeometry(QtCore.QRect(20, 60, 41, 31))
        self.inputDays.setAlignment(QtCore.Qt.AlignCenter)
        self.inputDays.setObjectName("inputDays")
        self.inputDays.setValidator(QIntValidator())
        self.inputDays.setMaxLength(3)
        self.inputDays_label = QtWidgets.QLabel(self.Deliverablesframe)
        self.inputDays_label.setGeometry(QtCore.QRect(70, 60, 101, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(False)
        font.setWeight(50)
        self.inputDays_label.setFont(font)
        self.inputDays_label.setObjectName("inputDays_label")
        self.deliverables_label = QtWidgets.QLabel(self.page)
        self.deliverables_label.setGeometry(QtCore.QRect(40, 482, 151, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.deliverables_label.setFont(font)
        self.deliverables_label.setObjectName("deliverables_label")
        self.ARMcheckBox = QtWidgets.QCheckBox(self.page)
        self.ARMcheckBox.setGeometry(QtCore.QRect(57, 552, 141, 21))
        self.ARMcheckBox.setChecked(True)
        self.ARMcheckBox.setObjectName("ARMcheckBox")
        self.ERBcheckBox = QtWidgets.QCheckBox(self.page)
        self.ERBcheckBox.setGeometry(QtCore.QRect(57, 522, 61, 21))
        self.ERBcheckBox.setChecked(True)
        self.ERBcheckBox.setObjectName("ERBcheckBox")
        self.markings_label = QtWidgets.QLabel(self.page)
        self.markings_label.setGeometry(QtCore.QRect(478, 154, 181, 51))
        font = QtGui.QFont()
        font.setPointSize(12)
        font.setBold(True)
        font.setWeight(75)
        self.markings_label.setFont(font)
        self.markings_label.setObjectName("markings_label")
        self.Classframe = QtWidgets.QFrame(self.page)
        self.Classframe.setGeometry(QtCore.QRect(656, 158, 201, 51))
        self.Classframe.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.Classframe.setFrameShadow(QtWidgets.QFrame.Raised)
        self.Classframe.setObjectName("Classframe")
        self.definedButton = QtWidgets.QRadioButton(self.Classframe)
        self.definedButton.setGeometry(QtCore.QRect(10, 26, 191, 25))
        self.definedButton.setChecked(False)
        self.definedButton.setObjectName("definedButton")
        self.unclassifiedButton = QtWidgets.QRadioButton(self.Classframe)
        self.unclassifiedButton.setGeometry(QtCore.QRect(10, 3, 131, 25))
        self.unclassifiedButton.setChecked(True)
        self.unclassifiedButton.setObjectName("unclassifiedButton")
        self.diagramDeleteButton = QtWidgets.QPushButton(self.page)
        self.diagramDeleteButton.setGeometry(QtCore.QRect(830, 238, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.diagramDeleteButton.setFont(font)
        self.diagramDeleteButton.setObjectName("diagramDeleteButton")
        self.listSUT_diagrams = QtWidgets.QListWidget(self.page)
        self.listSUT_diagrams.setGeometry(QtCore.QRect(476, 237, 350, 80))
        self.listSUT_diagrams.setDragEnabled(False)
        self.listSUT_diagrams.setDragDropOverwriteMode(False)
        self.listSUT_diagrams.setDragDropMode(QtWidgets.QAbstractItemView.NoDragDrop)
        self.listSUT_diagrams.setDefaultDropAction(QtCore.Qt.IgnoreAction)
        self.listSUT_diagrams.setObjectName("listSUT_diagrams")
        self.listSUT_diagrams_label = QtWidgets.QLabel(self.page)
        self.listSUT_diagrams_label.setGeometry(QtCore.QRect(478, 218, 191, 20))
        font = QtGui.QFont()
        font.setPointSize(12)
        font.setBold(True)
        font.setWeight(75)
        self.listSUT_diagrams_label.setFont(font)
        self.listSUT_diagrams_label.setObjectName("listSUT_diagrams_label")
        self.diagramUpButton = QtWidgets.QPushButton(self.page)
        self.diagramUpButton.setGeometry(QtCore.QRect(830, 267, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.diagramUpButton.setFont(font)
        self.diagramUpButton.setCursor(QtGui.QCursor(QtCore.Qt.ArrowCursor))
        self.diagramUpButton.setObjectName("diagramUpButton")
        self.diagramDownButton = QtWidgets.QPushButton(self.page)
        self.diagramDownButton.setGeometry(QtCore.QRect(830, 296, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.diagramDownButton.setFont(font)
        self.diagramDownButton.setObjectName("diagramDownButton")
        self.strengthsUpButton = QtWidgets.QPushButton(self.page)
        self.strengthsUpButton.setGeometry(QtCore.QRect(1290, 169, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.strengthsUpButton.setFont(font)
        self.strengthsUpButton.setCursor(QtGui.QCursor(QtCore.Qt.ArrowCursor))
        self.strengthsUpButton.setObjectName("strengthsUpButton")
        self.strengthsDownButton = QtWidgets.QPushButton(self.page)
        self.strengthsDownButton.setGeometry(QtCore.QRect(1290, 198, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.strengthsDownButton.setFont(font)
        self.strengthsDownButton.setObjectName("strengthsDownButton")
        self.weaknessesDownButton = QtWidgets.QPushButton(self.page)
        self.weaknessesDownButton.setGeometry(QtCore.QRect(1290, 418, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.weaknessesDownButton.setFont(font)
        self.weaknessesDownButton.setObjectName("weaknessesDownButton")
        self.weaknessesUpButton = QtWidgets.QPushButton(self.page)
        self.weaknessesUpButton.setGeometry(QtCore.QRect(1290, 389, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.weaknessesUpButton.setFont(font)
        self.weaknessesUpButton.setCursor(QtGui.QCursor(QtCore.Qt.ArrowCursor))
        self.weaknessesUpButton.setObjectName("weaknessesUpButton")
        self.mitigationsDownButton = QtWidgets.QPushButton(self.page)
        self.mitigationsDownButton.setGeometry(QtCore.QRect(1290, 638, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.mitigationsDownButton.setFont(font)
        self.mitigationsDownButton.setObjectName("mitigationsDownButton")
        self.mitigationsUpButton = QtWidgets.QPushButton(self.page)
        self.mitigationsUpButton.setGeometry(QtCore.QRect(1290, 609, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.mitigationsUpButton.setFont(font)
        self.mitigationsUpButton.setCursor(QtGui.QCursor(QtCore.Qt.ArrowCursor))
        self.mitigationsUpButton.setObjectName("mitigationsUpButton")
        self.reportcheckBox = QtWidgets.QCheckBox(self.page)
        self.reportcheckBox.setGeometry(QtCore.QRect(57, 582, 121, 21))
        self.reportcheckBox.setChecked(True)
        self.reportcheckBox.setObjectName("reportcheckBox")
        self.strengthsDeleteButton.raise_()
        self.strengths_label.raise_()
        self.listStrengths.raise_()
        self.strengthsAddButton.raise_()
        self.sutPreview_label.raise_()
        self.sutPreview_image.raise_()
        self.SUT_label.raise_()
        self.listWeaknesses.raise_()
        self.weaknesses_label.raise_()
        self.weaknessesAddButton.raise_()
        self.weaknessesDeleteButton.raise_()
        self.mitigationsDeleteButton.raise_()
        self.mitigationsAddButton.raise_()
        self.listMitigations.raise_()
        self.mitigations_label.raise_()
        self.inputStrength.raise_()
        self.inputStrength_label.raise_()
        self.inputWeakness_label.raise_()
        self.inputWeakness.raise_()
        self.inputMitigation_label.raise_()
        self.inputMitigation.raise_()
        self.saveButton.raise_()
        self.diagramSUT_label.raise_()
        self.figureBrowseButton.raise_()
        self.inputScope_label.raise_()
        self.listScope.raise_()
        self.scope_label.raise_()
        self.scopeAddButton.raise_()
        self.scopeDeleteButton.raise_()
        self.inputScope.raise_()
        self.Scopeframe.raise_()
        self.Deliverablesframe.raise_()
        self.deliverables_label.raise_()
        self.ARMcheckBox.raise_()
        self.ERBcheckBox.raise_()
        self.markings_label.raise_()
        self.Classframe.raise_()
        self.diagramDeleteButton.raise_()
        self.listSUT_diagrams.raise_()
        self.listSUT_diagrams_label.raise_()
        self.diagramUpButton.raise_()
        self.diagramDownButton.raise_()
        self.strengthsUpButton.raise_()
        self.strengthsDownButton.raise_()
        self.weaknessesDownButton.raise_()
        self.weaknessesUpButton.raise_()
        self.mitigationsDownButton.raise_()
        self.mitigationsUpButton.raise_()
        self.reportcheckBox.raise_()
        self.stackedWidget0.addWidget(self.page)
        MainWindow.setCentralWidget(self.centralwidget)
        self.menubar = QtWidgets.QMenuBar(MainWindow)
        self.menubar.setGeometry(QtCore.QRect(0, 0, 1345, 24))
        self.menubar.setObjectName("menubar")
        MainWindow.setMenuBar(self.menubar)
        self.statusbar = QtWidgets.QStatusBar(MainWindow)
        self.statusbar.setObjectName("statusbar")
        MainWindow.setStatusBar(self.statusbar)

        self.retranslateUi(MainWindow)
        self.stackedWidget0.setCurrentIndex(0)
        QtCore.QMetaObject.connectSlotsByName(MainWindow)
        MainWindow.setTabOrder(self.leadName, self.leadTitle)
        MainWindow.setTabOrder(self.leadTitle, self.leadOrg)
        MainWindow.setTabOrder(self.leadOrg, self.leadNIPR)
        MainWindow.setTabOrder(self.leadNIPR, self.leadSIPR)
        MainWindow.setTabOrder(self.leadSIPR, self.leadOffice)
        MainWindow.setTabOrder(self.leadOffice, self.leadMobile)
        MainWindow.setTabOrder(self.leadMobile, self.office_comboBox)
        MainWindow.setTabOrder(self.office_comboBox, self.eventName)
        MainWindow.setTabOrder(self.eventName, self.startDate)
        MainWindow.setTabOrder(self.startDate, self.endDate)
        MainWindow.setTabOrder(self.endDate, self.cvpaButton)
        MainWindow.setTabOrder(self.cvpaButton, self.pmrButton)
        MainWindow.setTabOrder(self.pmrButton, self.lightButton)
        MainWindow.setTabOrder(self.lightButton, self.darkButton)
        MainWindow.setTabOrder(self.darkButton, self.systemButton)
        MainWindow.setTabOrder(self.systemButton, self.unclassButton)
        MainWindow.setTabOrder(self.unclassButton, self.cuiButton)
        MainWindow.setTabOrder(self.cuiButton, self.fouoButton)
        MainWindow.setTabOrder(self.fouoButton, self.secretButton)
        MainWindow.setTabOrder(self.secretButton, self.noforncheckBox)
        MainWindow.setTabOrder(self.noforncheckBox, self.topsecretButton)
        MainWindow.setTabOrder(self.topsecretButton, self.scicheckBox)
        MainWindow.setTabOrder(self.scicheckBox, self.draftcheckBox)
        MainWindow.setTabOrder(self.draftcheckBox, self.listFolders)
        MainWindow.setTabOrder(self.listFolders, self.nextButton)
        MainWindow.setTabOrder(self.nextButton, self.listFindings)
        MainWindow.setTabOrder(self.listFindings, self.findingAddButton)
        MainWindow.setTabOrder(self.findingAddButton, self.findingUpButton)
        MainWindow.setTabOrder(self.findingUpButton, self.findingDownButton)
        MainWindow.setTabOrder(self.findingDownButton, self.findingDeleteButton)
        MainWindow.setTabOrder(self.findingDeleteButton, self.findingName)
        MainWindow.setTabOrder(self.findingName, self.findingHosts)
        MainWindow.setTabOrder(self.findingHosts, self.issuesText)
        MainWindow.setTabOrder(self.issuesText, self.insiderButton)
        MainWindow.setTabOrder(self.insiderButton, self.nearsiderButton)
        MainWindow.setTabOrder(self.nearsiderButton, self.outsiderButton)
        MainWindow.setTabOrder(self.outsiderButton, self.mitigationText)
        MainWindow.setTabOrder(self.mitigationText, self.mitigationcheckBox)
        MainWindow.setTabOrder(self.mitigationcheckBox, self.updateDescButton)
        MainWindow.setTabOrder(self.updateDescButton, self.listScreenshots)
        MainWindow.setTabOrder(self.listScreenshots, self.screenshotAddButton)
        MainWindow.setTabOrder(self.screenshotAddButton, self.screenshotUpButton)
        MainWindow.setTabOrder(self.screenshotUpButton, self.screenshotDownButton)
        MainWindow.setTabOrder(self.screenshotDownButton, self.screenshotDeleteButton)
        MainWindow.setTabOrder(self.screenshotDeleteButton, self.gobackButton)
        MainWindow.setTabOrder(self.gobackButton, self.quitButton)
        MainWindow.setTabOrder(self.quitButton, self.pptxButton)
        MainWindow.setTabOrder(self.pptxButton, self.inputScope)
        MainWindow.setTabOrder(self.inputScope, self.scopeAddButton)
        MainWindow.setTabOrder(self.scopeAddButton, self.bulletButton)
        MainWindow.setTabOrder(self.bulletButton, self.subbulletButton)
        MainWindow.setTabOrder(self.subbulletButton, self.listScope)
        MainWindow.setTabOrder(self.listScope, self.scopeDeleteButton)
        MainWindow.setTabOrder(self.scopeDeleteButton, self.ERBcheckBox)
        MainWindow.setTabOrder(self.ERBcheckBox, self.ARMcheckBox)
        MainWindow.setTabOrder(self.ARMcheckBox, self.reportcheckBox)
        MainWindow.setTabOrder(self.reportcheckBox, self.memoButton)
        MainWindow.setTabOrder(self.memoButton, self.reportButton)
        MainWindow.setTabOrder(self.reportButton, self.inputDays)
        MainWindow.setTabOrder(self.inputDays, self.figureBrowseButton)
        MainWindow.setTabOrder(self.figureBrowseButton, self.unclassifiedButton)
        MainWindow.setTabOrder(self.unclassifiedButton, self.definedButton)
        MainWindow.setTabOrder(self.definedButton, self.listSUT_diagrams)
        MainWindow.setTabOrder(self.listSUT_diagrams, self.diagramDeleteButton)
        MainWindow.setTabOrder(self.diagramDeleteButton, self.diagramUpButton)
        MainWindow.setTabOrder(self.diagramUpButton, self.diagramDownButton)
        MainWindow.setTabOrder(self.diagramDownButton, self.inputStrength)
        MainWindow.setTabOrder(self.inputStrength, self.strengthsAddButton)
        MainWindow.setTabOrder(self.strengthsAddButton, self.listStrengths)
        MainWindow.setTabOrder(self.listStrengths, self.strengthsDeleteButton)
        MainWindow.setTabOrder(self.strengthsDeleteButton, self.strengthsUpButton)
        MainWindow.setTabOrder(self.strengthsUpButton, self.strengthsDownButton)
        MainWindow.setTabOrder(self.strengthsDownButton, self.inputWeakness)
        MainWindow.setTabOrder(self.inputWeakness, self.weaknessesAddButton)
        MainWindow.setTabOrder(self.weaknessesAddButton, self.listWeaknesses)
        MainWindow.setTabOrder(self.listWeaknesses, self.weaknessesDeleteButton)
        MainWindow.setTabOrder(self.weaknessesDeleteButton, self.weaknessesUpButton)
        MainWindow.setTabOrder(self.weaknessesUpButton, self.weaknessesDownButton)
        MainWindow.setTabOrder(self.weaknessesDownButton, self.inputMitigation)
        MainWindow.setTabOrder(self.inputMitigation, self.mitigationsAddButton)
        MainWindow.setTabOrder(self.mitigationsAddButton, self.listMitigations)
        MainWindow.setTabOrder(self.listMitigations, self.mitigationsDeleteButton)
        MainWindow.setTabOrder(self.mitigationsDeleteButton, self.mitigationsUpButton)
        MainWindow.setTabOrder(self.mitigationsUpButton, self.mitigationsDownButton)
        MainWindow.setTabOrder(self.mitigationsDownButton, self.saveButton)


        class Event:
            def __init__(self, lead_name, lead_title, lead_org, lead_nipr, lead_sipr, lead_office, lead_mobile, office_symbol, event_name, event_type, start_date, end_date, classification, designation, draft, mode):
                self.lead_name = lead_name
                self.lead_title = lead_title
                self.lead_org = lead_org
                self.lead_nipr = lead_nipr
                self.lead_sipr = lead_sipr
                self.lead_office = lead_office
                self.lead_mobile = lead_mobile
                self.office_symbol = office_symbol
                self.event_name = event_name
                self.event_type = event_type
                self.start_date = start_date
                self.end_date = end_date
                self.classification = classification
                self.designation = designation
                self.draft = draft
                self.mode = mode

            def write_file(self):
                #CREATE NEW XML FILE
                xml_file = cwd + '/event.xml'
                def indent(elem, level=0):
                    i = "\n" + level*"    "
                    if len(elem):
                        if not elem.text or not elem.text.strip():
                            elem.text = i + "    "
                        if not elem.tail or not elem.tail.strip():
                            elem.tail = i
                        for elem in elem:
                            indent(elem, level+1)
                        if not elem.tail or not elem.tail.strip():
                            elem.tail = i
                    else:
                        if level and (not elem.tail or not elem.tail.strip()):
                            elem.tail = i
                #CREATE FILE STRUCTURE
                root = ET.Element('data')
                #CREATE EVENT
                xml_event = ET.SubElement(root, 'event')
                xml_event.set('uid', str(0))
                xml_lead_name = ET.SubElement(xml_event, 'lead_name')
                xml_lead_name.text = self.lead_name
                xml_lead_title = ET.SubElement(xml_event, 'lead_title')
                xml_lead_title.text = self.lead_title
                xml_lead_org = ET.SubElement(xml_event, 'lead_org')
                xml_lead_org.text = self.lead_org
                xml_lead_nipr = ET.SubElement(xml_event, 'lead_nipr')
                xml_lead_nipr.text = self.lead_nipr
                xml_lead_sipr = ET.SubElement(xml_event, 'lead_sipr')
                xml_lead_sipr.text = self.lead_sipr
                xml_lead_office = ET.SubElement(xml_event, 'lead_office')
                xml_lead_office.text = self.lead_office
                xml_lead_mobile = ET.SubElement(xml_event, 'lead_mobile')
                xml_lead_mobile.text = self.lead_mobile
                xml_office_symbol = ET.SubElement(xml_event, 'office_symbol')
                xml_office_symbol.text = self.office_symbol
                xml_event_name = ET.SubElement(xml_event, 'event_name')
                xml_event_name.text = self.event_name
                xml_event_type = ET.SubElement(xml_event, 'event_type')
                xml_event_type.text = self.event_type
                xml_start_date = ET.SubElement(xml_event, 'start_date')
                xml_start_date.text = self.start_date
                xml_end_date = ET.SubElement(xml_event, 'end_date')
                xml_end_date.text = self.end_date
                xml_classification = ET.SubElement(xml_event, 'classification')
                xml_classification.text = self.classification
                xml_designation = ET.SubElement(xml_event, 'designation')
                xml_designation.text = self.designation
                xml_draft = ET.SubElement(xml_event, 'draft')
                xml_draft.text = self.draft
                xml_mode = ET.SubElement(xml_event, 'mode')
                xml_mode.text = self.mode
                #WRITING XML
                indent(root)
                tree = ET.ElementTree(root)
                tree.write(xml_file, encoding='utf-8', xml_declaration=True)


        class Finding:
            def __init__(self, folder, active, rank, title, hosts, issues, posture, mitigation, include_mitigation, screenshots=None):
                self.folder = folder
                self.active = active
                self.rank = rank
                self.title = title
                self.hosts = hosts
                self.issues = issues
                self.posture = posture
                self.mitigation = mitigation
                self.include_mitigation = include_mitigation
                self.screenshots = screenshots

            def insert_new_finding(self, uid):
                def indent(elem, level=0):
                    i = "\n" + level*"    "
                    if len(elem):
                        if not elem.text or not elem.text.strip():
                            elem.text = i + "    "
                        if not elem.tail or not elem.tail.strip():
                            elem.tail = i
                        for elem in elem:
                            indent(elem, level+1)
                        if not elem.tail or not elem.tail.strip():
                            elem.tail = i
                    else:
                        if level and (not elem.tail or not elem.tail.strip()):
                            elem.tail = i
                #UPDATE XML FILE
                tree = ET.parse(current_erb)
                root = tree.getroot()
                xml_finding = ET.SubElement(root, 'finding')
                xml_finding.set('uid', uid)
                xml_folder = ET.SubElement(xml_finding, 'folder')
                xml_folder.text = str(self.folder)
                xml_active = ET.SubElement(xml_finding, 'active')
                xml_active.text = str(self.active)
                xml_rank = ET.SubElement(xml_finding, 'rank')
                xml_rank.text = str(self.rank)
                xml_title = ET.SubElement(xml_finding, 'title')
                xml_title.text = str(self.title)
                xml_hosts = ET.SubElement(xml_finding, 'hosts')
                xml_hosts.text = str(self.hosts)
                xml_issues = ET.SubElement(xml_finding, 'issues')
                xml_issues.text = str(self.issues)
                xml_posture = ET.SubElement(xml_finding, 'posture')
                xml_posture.text = str(self.posture)
                xml_mitigation = ET.SubElement(xml_finding, 'mitigation')
                xml_mitigation.text = str(self.mitigation)
                xml_include_mitigation = ET.SubElement(xml_finding, 'include_mitigation')
                xml_include_mitigation.text = str(self.include_mitigation)
                xml_screenshots = ET.SubElement(xml_finding, 'screenshots')
                xml_screenshots.text = str(self.screenshots)
                #WRITING XML
                indent(root)
                tree.write(current_erb, encoding='utf-8', xml_declaration=True)

            def update_attributes(self, finding_num, title, hosts, issues, posture, mitigation, include_mitigation):
                self.title = title
                self.hosts = hosts
                self.issues = issues
                self.posture = posture
                self.mitigation = mitigation
                self.include_mitigation = include_mitigation
                #UPDATE FINDING IN XML
                tree = ET.parse(current_erb)
                root = tree.getroot()
                for finding in root.findall('finding'):
                    uid = finding.get('uid')
                    if uid == finding_num:
                        xml_title = finding.find('title')
                        xml_hosts = finding.find('hosts')
                        xml_issues = finding.find('issues')
                        xml_posture = finding.find('posture')
                        xml_mitigation = finding.find('mitigation')
                        xml_include_mitigation = finding.find('include_mitigation')
                        xml_title.text = self.title
                        xml_hosts.text = self.hosts
                        xml_issues.text = self.issues
                        xml_posture.text = self.posture
                        xml_mitigation.text = self.mitigation
                        xml_include_mitigation.text = self.include_mitigation
                tree.write(current_erb)

            @classmethod
            def modify_rank(self):
                #MODIFY RANK ORDER IN XML
                tree = ET.parse(current_erb)
                root = tree.getroot()
                for finding in root.findall('finding'):
                    uid = finding.get('uid')
                    active = finding.find('active').text
                    if active == '1':
                        rank = finding.find('rank')
                        value = order_list.index(uid)
                        rank.text = str(value)
                tree.write(current_erb)

            @classmethod
            def deactivate(self, finding_num):
                #DE-ACTIVATE FINDING IN XML
                tree = ET.parse(current_erb)
                root = tree.getroot()
                for finding in root.findall('finding'):
                    uid = finding.get('uid')
                    if uid == order_list[finding_num]:
                        active = finding.find('active')
                        rank = finding.find('rank')
                        active.text = '0'
                        rank.text = 'x'
                tree.write(current_erb)
                order_list.remove(order_list[finding_num])

            def modify_folder(self, finding_num):
                #MODIFY SCREENSHOTS FOLDER IN XML
                tree = ET.parse(current_erb)
                root = tree.getroot()
                for finding in root.findall('finding'):
                    uid = finding.get('uid')
                    if uid == order_list[finding_num]:
                        screenshots = finding.find('folder')
                        screenshots.text = str(self.folder)
                tree.write(current_erb)

            def modify_screenshots(self, finding_num):
                #MODIFY SCREENSHOTS ORDER IN XML
                tree = ET.parse(current_erb)
                root = tree.getroot()
                for finding in root.findall('finding'):
                    uid = finding.get('uid')
                    if uid == order_list[finding_num]:
                        screenshots = finding.find('screenshots')
                        screenshots.text = str(self.screenshots)
                tree.write(current_erb)


        #CLEAR ALL FIELDS ON PAGE 2
        def clear_all():
            self.listFindings.clear()
            self.findingName.setText("")
            self.findingHosts.setText("")
            self.issuesText.clear()
            if self.nearsiderButton.isChecked() == False:
                self.nearsiderButton.toggle()
            self.mitigationText.clear()
            if self.mitigationcheckBox.isChecked() == True:
                self.mitigationcheckBox.toggle()
            self.listScreenshots.clear()
            self.imagePreview_label.setText("NO PREVIEW")


        #DISPLAY IMAGES ON "SCREENSHOT PREVIEW"
        def preview_image():
            if self.listScreenshots.selectedItems() != []:
                item1 = self.listFindings.selectedIndexes()[0]
                finding_num = item1.row()
                current_finding = order_list[finding_num]
                item2 = self.listScreenshots.selectedIndexes()[0]
                image_num = item2.row()
                image_name = findings[current_finding].screenshots[image_num]
                image_file = findings[str(order_list[finding_num])].folder + image_name
                self.imagePreview_label.setText("")
                self.imagePreview_label.setPixmap(QtGui.QPixmap(image_file))
                self.imagePreview_label.setScaledContents(True)
            else:
                self.imagePreview_label.setText("NO PREVIEW")


        #POPULATE SCREENSHOTS LIST
        def set_screenshots():
            self.listScreenshots.clear()
            item = self.listFindings.selectedIndexes()[0]
            index = item.row()
            current_finding = order_list[index]
            if findings[current_finding].screenshots != None:
                for i in findings[current_finding].screenshots:
                    self.listScreenshots.addItem(i)
                if findings[current_finding].screenshots != [] and self.listScreenshots.selectedItems() == []:
                    self.listScreenshots.item(0).setSelected(True)
                    self.listScreenshots.setCurrentRow(0)


        #POPULATE FINDINGS FIELDS
        def set_fields():
            if self.listFindings.selectedItems() != []:
                item = self.listFindings.selectedIndexes()[0]
                index = item.row()
                current_finding = order_list[index]
                self.findingName.setText(findings[current_finding].title)
                self.findingHosts.setText(findings[current_finding].hosts)
                self.issuesText.clear()
                self.issuesText.insertPlainText(findings[current_finding].issues)
                if findings[current_finding].posture == 'INSIDER':
                    if self.insiderButton.isChecked() == False:
                        self.insiderButton.toggle()
                if findings[current_finding].posture == 'NEARSIDER':
                    if self.nearsiderButton.isChecked() == False:
                        self.nearsiderButton.toggle()
                if findings[current_finding].posture == 'OUTSIDER':
                    if self.outsiderButton.isChecked() == False:
                        self.outsiderButton.toggle()
                self.mitigationText.clear()
                self.mitigationText.insertPlainText(findings[current_finding].mitigation)
                if findings[current_finding].include_mitigation == 'yes':
                    if self.mitigationcheckBox.isChecked() == False:
                        self.mitigationcheckBox.toggle()
                if findings[current_finding].include_mitigation == 'no':
                    if self.mitigationcheckBox.isChecked() == True:
                        self.mitigationcheckBox.toggle()
                set_screenshots()


        #UPDATE FINDING XML
        def update_finding():
            if self.listFindings.selectedItems() != []:
                item = self.listFindings.selectedIndexes()[0]
                index = item.row()
                current_finding = order_list[index]
                title = self.findingName.text()
                hosts = self.findingHosts.text()
                issues = self.issuesText.toPlainText()
                if self.insiderButton.isChecked() == True:
                    posture = 'INSIDER'
                if self.nearsiderButton.isChecked() == True:
                    posture = 'NEARSIDER'
                if self.outsiderButton.isChecked() == True:
                    posture = 'OUTSIDER'
                mitigation = self.mitigationText.toPlainText()
                if self.mitigationcheckBox.isChecked() == True:
                    include_mitigation = 'yes'
                else:
                    include_mitigation = 'no'
                #UPDATE FINDING
                findings[str(current_finding)].update_attributes(current_finding, title, hosts, issues, posture, mitigation, include_mitigation)
                self.listFindings.clear()
                for i in order_list:
                    self.listFindings.addItem(findings[i].title)
                self.listFindings.item(index).setSelected(True)
                self.listFindings.setCurrentRow(index)
                #MESSAGE BOX!
                msg = QMessageBox()
                msg.setIcon(QMessageBox.Information)
                msg.setWindowTitle("Update Finding")
                msg.setText("Finding was successfully updated!")
                x = msg.exec_()


        #GET AVAILABLE DRADIS/FRIC FILES/FOLDERS FROM "DESKTOP" + CURRENT ERB
        def get_folders():
            global cwd, current_erb
            if os.path.isfile(current_erb):
                self.listFolders.addItem('Existing ERB')
            self.listFolders.addItem('Create your own ERB')
            os.chdir(desktop_dir)
            for file in glob.glob("dradis-export*.zip"):
                self.listFolders.addItem(str(file))
            raw_folders = next(os.walk(desktop_dir))[1]
            prefix = "fric_export_"
            for i in raw_folders:
                if prefix in i:
                    self.listFolders.addItem(str(i))
            os.chdir(cwd)
        get_folders()


        #GET FINDINGS FROM DRADIS/FRIC FILES/FOLDERS
        def get_findings():
            global order_list
            tree = ET.parse(current_erb)
            root = tree.getroot()
            num_findings = len(root)
            findings.clear()
            self.listFindings.clear()
            uids = []
            ranks = []
            del order_list[:]
            for finding in root.findall('finding'):
                uid = finding.get('uid')
                folder = finding.find('folder').text
                active = finding.find('active').text
                rank = finding.find('rank').text
                title = finding.find('title').text
                hosts = finding.find('hosts').text
                issues = finding.find('issues').text
                posture = finding.find('posture').text
                mitigation = finding.find('mitigation').text
                include_mitigation = finding.find('include_mitigation').text
                tmp = finding.find('screenshots').text
                if tmp != None:
                    tmp = tmp.split('[')
                    tmp = tmp[1].split(']')
                    tmp = tmp[0]
                    tmp = tmp.split(', ')
                    for k in range(len(tmp)):
                        tmp2 = tmp[k].split("'")
                        if len(tmp2) > 1:
                            tmp[k] = tmp2[1]
                            screenshots = tmp
                        else:
                            screenshots = None
                else:
                    screenshots = None
                if active == '0':
                    rank = 'x'
                else:
                    findings[uid] = Finding(folder, active, rank, title, hosts, issues, posture, mitigation, include_mitigation, screenshots)
                    uids.append(uid)
                    ranks.append(rank)
            int_ranks = [int(x) for x in ranks]
            order_list = [x for _,x in sorted(zip(int_ranks,uids))]
            #POPULATE FINDINGS LIST
            for i in order_list:
                self.listFindings.addItem(findings[i].title)


        #CHECK IF ITEM IS SELECTED FROM DRADIS/FRIC FILES/FOLDERS
        def on_selection_changed():
            clear_all()
            global data_folder, data_source
            data_source = self.listFolders.currentItem().text()
            if self.listFolders.selectedItems():
                self.nextButton.setEnabled(True)
                if 'Existing ERB' in data_source:
                    data_folder = cwd + '/erb/'
                elif 'Create your own ERB' in data_source:
                    data_folder = cwd + '/erb/'
                else:
                    data_folder = desktop_dir + '/' + data_source


        #CHECK IF EVENT XML FILE EXISTS
        xml_file = cwd + '/event.xml'
        if os.path.isfile(xml_file):
            tree = ET.parse(xml_file)
            root = tree.getroot()
            for event in root.findall('event'):
                uid = event.get('uid')
                lead_name = event.find('lead_name').text
                lead_title = event.find('lead_title').text
                lead_org = event.find('lead_org').text
                lead_nipr = event.find('lead_nipr').text
                lead_sipr = event.find('lead_sipr').text
                lead_office = event.find('lead_office').text
                lead_mobile = event.find('lead_mobile').text
                office_symbol = event.find('office_symbol').text
                event_name = event.find('event_name').text
                event_type = event.find('event_type').text
                start_date = event.find('start_date').text
                end_date = event.find('end_date').text
                classification = event.find('classification').text
                designation = event.find('designation').text
                draft = event.find('draft').text
                mode = event.find('mode').text
                
                self.leadName.setText(lead_name)
                self.leadTitle.setText(lead_title)
                self.leadOrg.setText(lead_org)
                self.leadNIPR.setText(lead_nipr)
                self.leadSIPR.setText(lead_sipr)
                self.leadOffice.setText(lead_office)
                self.leadMobile.setText(lead_mobile)
                self.office_comboBox.setCurrentText(office_symbol)
                self.eventName.setText(event_name)
                if event_type == 'CVPA':
                    self.cvpaButton.setChecked(True)
                else:
                    self.pmrButton.setChecked(True)
                    self.systemButton.hide()
                    self.dradisfricData_label.hide()
                    self.availFilesFolders_label.hide()
                    self.listFolders.hide()
                    self.nextButton.setText("PPTX")
                    self.nextButton.setEnabled(True)
                start_date = start_date.split('/')
                self.startDate.setDate(QDate(int(start_date[2]), int(start_date[0]), int(start_date[1])))
                end_date = end_date.split('/')
                self.endDate.setDate(QDate(int(end_date[2]), int(end_date[0]), int(end_date[1])))
                self.cuiButton.setEnabled(False)
                self.fouoButton.setEnabled(False)
                if classification == 'UNCLASSIFIED':
                    self.unclassButton.setChecked(True)
                    self.cuiButton.setEnabled(True)
                    self.fouoButton.setEnabled(True)
                    if designation == 'CUI':
                        self.cuiButton.setChecked(True)
                    elif designation == 'FOUO':
                        self.fouoButton.setChecked(True)
                    else:
                        self.fouoButton.setChecked(False)
                if classification == 'SECRET':
                    self.secretButton.setChecked(True)
                    self.noforncheckBox.setEnabled(True)
                    if designation == 'NOFORN':
                        self.noforncheckBox.setChecked(True)
                if classification == 'TOP SECRET':
                    self.topsecretButton.setChecked(True)
                    self.scicheckBox.setEnabled(True)
                    if designation == 'SCI':
                        self.scicheckBox.setChecked(True)
                if draft == 'YES':
                    self.draftcheckBox.setChecked(True)
                else:
                    self.draftcheckBox.setChecked(False)
                if mode == 'LIGHT':
                    if self.lightButton.isChecked() == False:
                        self.lightButton.toggle()
                if mode == 'DARK':
                    if self.darkButton.isChecked() == False:
                        self.darkButton.toggle()
        else:
            #HANDLE DATES
            #TODAY'S DATE
            now = QDate.currentDate()
            #TODAY'S WEEK DAY - GETTING INDEX/OFFSET
            week_days = ['Mon', 'Tue', 'Wed', 'Thu', 'Fri', 'Sat', 'Sun']
            datetime = QDateTime.currentDateTime()
            datetime = datetime.toString()
            week_day = datetime[0:3]
            todays_index = week_days.index(week_day) #Mon - 0, Tue - 1, ..., Sun - 6
            #DETERMINE START DAY
            start_day = now.addDays(-todays_index)
            #DETERMINE END DAY
            end_day = start_day.addDays(4)
            #SETTING NEW DATES - START / END
            self.startDate.setDate(start_day)
            self.endDate.setDate(end_day)


        #DELIVERABLES - REPORT/MEMO
        def deliverables_changed():
            if self.reportcheckBox.isChecked() == True:
                self.memoButton.setEnabled(True)
                self.reportButton.setEnabled(True)
                self.inputDays.setEnabled(True)
                self.inputDays_label.setEnabled(True)
            else:
                self.memoButton.setEnabled(False)
                self.reportButton.setEnabled(False)
                self.inputDays.setEnabled(False)
                self.inputDays_label.setEnabled(False)


        #PARSE DATA FROM XML
        def parse_xml(data):
            data = data.replace('[', '')
            data = data.replace(']', '')
            data = data.replace("', '", "'!@#$% '")
            data = data.split('!@#$% ')
            aux = []
            for i in data:
                aux.append(i[1:-1])
            return aux


        #CHECK IF SUT INFO XML FILE EXISTS
        global scope_list, diagram_marking, diagrams_list, strengths_list, weaknesses_list, mitigations_list, deliverables
        xml_file = cwd + '/sut_info.xml'
        if os.path.isfile(xml_file):
            tree = ET.parse(xml_file)
            root = tree.getroot()
            for info in root.findall('sut'):
                uid = info.get('uid')
                scope_list = parse_xml(info.find('scope').text)
                if len(scope_list) == 1 and scope_list[0] == '':
                    scope_list.clear()
                else:
                    for i in scope_list:
                        self.listScope.addItem(i)
                diagram_marking = info.find('diagram').text
                if diagram_marking == 'UNCLASSIFIED':
                    self.unclassifiedButton.setChecked(True)
                else:
                    self.definedButton.setChecked(True)
                diagrams_list = parse_xml(info.find('diagrams_list').text)
                if len(diagrams_list) == 1 and diagrams_list[0] == '':
                    diagrams_list.clear()
                else:
                    if diagrams_list != 0:
                        for i in diagrams_list:
                            self.listSUT_diagrams.addItem(i)
                            diagrams_list = [str(self.listSUT_diagrams.item(i).text()) for i in range(self.listSUT_diagrams.count())]
                        self.listSUT_diagrams.item(0).setSelected(True)
                        self.listSUT_diagrams.setCurrentRow(0)
                        sut_file = cwd + '/' + diagrams_list[0]
                        if os.path.isfile(sut_file):
                            self.sutPreview_image.setText("")
                            self.sutPreview_image.setPixmap(QtGui.QPixmap(sut_file))
                            self.sutPreview_image.setScaledContents(True)
                        else:
                            self.sutPreview_image.setText("NO PREVIEW")
                strengths_list = parse_xml(info.find('strengths').text)
                if len(strengths_list) == 1 and strengths_list[0] == '':
                    strengths_list.clear()
                else:
                    for i in strengths_list:
                        self.listStrengths.addItem(i)
                weaknesses_list = parse_xml(info.find('weaknesses').text)
                if len(weaknesses_list) == 1 and weaknesses_list[0] == '':
                    weaknesses_list.clear()
                else:
                    for i in weaknesses_list:
                        self.listWeaknesses.addItem(i)
                mitigations_list = parse_xml(info.find('mitigations').text)
                if len(mitigations_list) == 1 and mitigations_list[0] == '':
                    mitigations_list.clear()
                else:
                    for i in mitigations_list:
                        self.listMitigations.addItem(i)
                deliverables = parse_xml(info.find('deliverables').text)
                if deliverables[0] == 'YES':
                    self.ERBcheckBox.setChecked(True)
                else:
                    self.ERBcheckBox.setChecked(False)
                if deliverables[1] == 'YES':
                    self.ARMcheckBox.setChecked(True)
                else:
                    self.ARMcheckBox.setChecked(False)
                if deliverables[2] == 'YES':
                    self.reportcheckBox.setChecked(True)
                else:
                    self.reportcheckBox.setChecked(False)
                deliverables_changed()
                if deliverables[3] == 'MEMO':
                    self.memoButton.setChecked(True)
                else:
                    self.reportButton.setChecked(True)
                self.inputDays.setText(deliverables[4])
        else:
            self.listScope.clear()
            self.unclassifiedButton.setChecked(True)
            self.sutPreview_image.setText("NO PREVIEW")
            self.listStrengths.clear()
            self.listWeaknesses.clear()
            self.listMitigations.clear()
            self.ERBcheckBox.setChecked(True)
            self.ARMcheckBox.setChecked(True)
            self.reportcheckBox.setChecked(True)
            self.reportButton.setChecked(True)
            self.inputDays.setText("90")


        #NIPR EMAIL CHANGE
        def nipr_changed():
            n_email = self.leadNIPR.text()
            new_email = n_email.split('@')
            s_email = new_email[0] + '@mail.smil.mil'
            self.leadSIPR.setText(s_email)


        #EVENT TYPE SELECTION
        def event_type():
            if self.cvpaButton.isChecked() == True:
                self.systemButton.show()
                self.dradisfricData_label.show()
                self.availFilesFolders_label.show()
                self.listFolders.show()
                self.nextButton.setText("Next")
                if self.listFolders.selectedItems():
                    self.nextButton.setEnabled(True)
                else:
                    self.nextButton.setEnabled(False)
            if self.pmrButton.isChecked() == True:
                self.systemButton.hide()
                self.dradisfricData_label.hide()
                self.availFilesFolders_label.hide()
                self.listFolders.hide()
                self.nextButton.setText("PPTX")
                self.nextButton.setEnabled(True)


        #CLASSIFICATION SELECTION CHANGE
        def class_changed():
            self.noforncheckBox.setChecked(False)
            self.scicheckBox.setChecked(False)
            if self.unclassButton.isChecked():
                self.cuiButton.setEnabled(True)
                self.fouoButton.setEnabled(True)
                self.noforncheckBox.setEnabled(False)
                self.scicheckBox.setEnabled(False)
            if self.secretButton.isChecked():
                self.cuiButton.setEnabled(False)
                self.fouoButton.setEnabled(False)
                self.cuiButton.setChecked(False)
                self.fouoButton.setChecked(False)
                self.noforncheckBox.setEnabled(True)
                self.scicheckBox.setEnabled(False)
            if self.topsecretButton.isChecked():
                self.cuiButton.setEnabled(False)
                self.fouoButton.setEnabled(False)
                self.cuiButton.setChecked(False)
                self.fouoButton.setChecked(False)
                self.noforncheckBox.setEnabled(False)
                self.scicheckBox.setEnabled(True)


	#ADD FINDING TO LIST
        def addFinding():
            global order_list
            uid = str(len(findings))
            folder = ''
            active = 1
            rank = str(len(findings))
            title = 'NEW FINDING'
            hosts = 'ADD HOST(S)'
            issues = 'ADD ISSUE(S)'
            posture = 'NEARSIDER'
            mitigation = 'NO MITIGATION FOUND FOR THIS FINDING...'
            include_mitigation = 'yes'
            screenshots = []
            findings[uid] = Finding(folder, active, rank, title, hosts, issues, posture, mitigation, include_mitigation, screenshots)
            findings[uid].insert_new_finding(uid)
            order_list.append(uid)
            self.listFindings.addItem(title)
            current_list = [str(self.listFindings.item(i).text()) for i in range(self.listFindings.count())]
            index = len(current_list)-1
            self.listFindings.item(index).setSelected(True)
            self.listFindings.setCurrentRow(index)


        #MOVE ONE FINDING UP THE LIST
        def findingsUp():
            global order_list
            currentRow = self.listFindings.currentRow()
            currentItem = self.listFindings.takeItem(currentRow)
            if currentRow > 0:
                tmp = order_list[currentRow-1]
                order_list[currentRow-1] = order_list[currentRow]
                order_list[currentRow] = tmp
            if currentRow != -1:
                self.listFindings.insertItem(currentRow - 1, currentItem)
                self.listFindings.setCurrentItem(currentItem)
                selected_finding = self.listFindings.currentItem().text()
                current_list = [str(self.listFindings.item(i).text()) for i in range(self.listFindings.count())]
                Finding.modify_rank()


        #MOVE ONE FINDING DOWN THE LIST
        def findingsDown():
            global order_list
            currentRow = self.listFindings.currentRow()
            currentItem = self.listFindings.takeItem(currentRow)
            if currentRow < len(order_list)-1:
                tmp = order_list[currentRow+1]
                order_list[currentRow+1] = order_list[currentRow]
                order_list[currentRow] = tmp
            if currentRow != -1:
                self.listFindings.insertItem(currentRow + 1, currentItem)
                self.listFindings.setCurrentItem(currentItem)
                selected_finding = self.listFindings.currentItem().text()
                current_list = [str(self.listFindings.item(i).text()) for i in range(self.listFindings.count())]
                Finding.modify_rank()


        #DELETE FINDING FROM LIST
        def delFinding():
            global order_list
            currentRow = self.listFindings.currentRow()
            item = self.listFindings.takeItem(self.listFindings.currentRow())
            item = None
            current_list = [str(self.listFindings.item(i).text()) for i in range(self.listFindings.count())]
            if currentRow != -1:
                Finding.deactivate(currentRow)
            if current_list == []:
                clear_all()


        #ADD SCREENSHOT TO LIST
        def addScreenshot():
            global f_folder, screenshots
            if self.listFindings.selectedItems() != []:
                currentRow = self.listFindings.currentRow()
                if f_folder == '':
                    f_folder = desktop_dir
                fname, _filter = QtWidgets.QFileDialog.getOpenFileName(None, "Import Image", f_folder, "*.png *.jpg *.bmp *.gif *.jpeg *.tif *.tiff")
                if fname != '':
                    aux = fname.rsplit('/', 1)
                    f_folder = aux[0]
                    new_image = fname.split('/')
                    new_image = new_image[-1]
                    findings[order_list[currentRow]].screenshots.append(new_image)
                    if findings[order_list[currentRow]].folder == '':
                        sf = cwd + '/erb/'
                        num_folders = len(next(os.walk(sf))[1])
                        img_folder = sf + str(num_folders)
                        msg = 'mkdir ' + img_folder
                        os.system(msg)
                        findings[order_list[currentRow]].folder = img_folder + '/'
                        findings[order_list[currentRow]].modify_folder(currentRow)
                    else:
                        img_folder = findings[order_list[currentRow]].folder
                    fname = fname.replace('(', '\(')
                    fname = fname.replace(')', '\)')
                    fname = fname.replace(' ', '\ ')
                    msg = 'cp ' + str(fname) + ' ' + str(img_folder)
                    os.system(msg)
                    findings[order_list[currentRow]].modify_screenshots(currentRow)
                    self.listScreenshots.addItem(new_image)
                    current_list = [str(self.listScreenshots.item(i).text()) for i in range(self.listScreenshots.count())]
                    index = len(current_list)-1
                    self.listScreenshots.item(index).setSelected(True)
                    self.listScreenshots.setCurrentRow(index)


        #MOVE ONE SCREENSHOT UP THE LIST
        def screenshotsUp():
            if self.listScreenshots.selectedItems() != []:
                currentRow = self.listScreenshots.currentRow()
                currentItem = self.listScreenshots.takeItem(currentRow)
                self.listScreenshots.insertItem(currentRow - 1, currentItem)
                self.listScreenshots.setCurrentItem(currentItem) 
                finding_item = self.listFindings.selectedIndexes()[0]
                finding_num = finding_item.row()
                findings[order_list[finding_num]].screenshots.clear()
                current_list = [str(self.listScreenshots.item(i).text()) for i in range(self.listScreenshots.count())]
                for i in current_list:
                    findings[order_list[finding_num]].screenshots.append(i)
                preview_image()
                findings[order_list[finding_num]].modify_screenshots(finding_num)


        #MOVE ONE SCREENSHOT DOWN THE LIST
        def screenshotsDown():
            if self.listScreenshots.selectedItems() != []:
                currentRow = self.listScreenshots.currentRow()
                currentItem = self.listScreenshots.takeItem(currentRow)
                self.listScreenshots.insertItem(currentRow + 1, currentItem)
                self.listScreenshots.setCurrentItem(currentItem)
                finding_item = self.listFindings.selectedIndexes()[0]
                finding_num = finding_item.row()
                findings[order_list[finding_num]].screenshots.clear()
                current_list = [str(self.listScreenshots.item(i).text()) for i in range(self.listScreenshots.count())]
                for i in current_list:
                    findings[order_list[finding_num]].screenshots.append(i)
                preview_image()
                findings[order_list[finding_num]].modify_screenshots(finding_num)


        #DELETE SCREENSHOT FROM LIST
        def delScreenshot():
            global screenshots
            if self.listScreenshots.selectedItems() != []:
                item = self.listScreenshots.takeItem(self.listScreenshots.currentRow())
                item = None
                finding_item = self.listFindings.selectedIndexes()[0]
                finding_num = finding_item.row()
                findings[order_list[finding_num]].screenshots.clear()
                current_list = [str(self.listScreenshots.item(i).text()) for i in range(self.listScreenshots.count())]
                for i in current_list:
                    findings[order_list[finding_num]].screenshots.append(i)
                findings[order_list[finding_num]].modify_screenshots(finding_num)


        #POWERPOINT PRESENTATION
        def create_pptx():
            global event, diagrams_list
            #OPEN/CREATE PRESENTATION
            i_file = cwd + '/ERB_Template.pptx'
            prs = Presentation(i_file)
            datetime = QDateTime.currentDateTime()
            datetime = datetime.toString()

            #OUTPUT FILE NAME
            prefix = ''
            if event.event_type == 'CVPA':
                t_event = event.event_name.lower()
                if 'cvi' in t_event:
                    prefix = 'CVI_'
                elif 'vof' in t_event:
                    prefix = 'VOF_'
                else:
                    prefix = 'CVPA_'
            if event.event_type == 'PMR':
                prefix = 'PMR_'
            #SLIDE DECK MODE - LIGHT / DARK
            if event.mode == 'DARK':
                slide_deck = [9, 10, 11, 12, 13, 14, 15, 16, 17]
                c_color = [255, 255, 255]
                pptx_file = prefix + 'ERB_dark_' + datetime + '.pptx'
                pptx_file = pptx_file.replace(":", "_")
                pptx_file = pptx_file.replace(" ", "_")
            else:
                slide_deck = [0, 1, 2, 3, 4, 5, 6, 7, 8]
                c_color = [0, 0, 0]
                pptx_file = prefix + 'ERB_white_' + datetime + '.pptx'
                pptx_file = pptx_file.replace(":", "_")
                pptx_file = pptx_file.replace(" ", "_")

            #LAYOUTS
            cover_slide_layout = prs.slide_layouts[slide_deck[0]] #COVER SLIDE
            cover_slide_cui_layout = prs.slide_layouts[slide_deck[1]] #COVER SLIDE (CUI)
            title_subtitle_slide_layout = prs.slide_layouts[slide_deck[2]] #TITLE & SUBTITLE
            title_paragraph_bullets_slide_layout = prs.slide_layouts[slide_deck[3]] #TITLE & PARAGRAPH BULLETS
            title_paragraph_no_bullets_slide_layout = prs.slide_layouts[slide_deck[4]] #TITLE & PARAGRAPH NO BULLETS
            blank_slide_layout = prs.slide_layouts[slide_deck[5]] #BLANK SLIDE WITH TITLE
            picture_slide_layout = prs.slide_layouts[slide_deck[6]] #PICTURE SLIDE
            pen_process_slide_layout = prs.slide_layouts[slide_deck[7]] #PENETRATION PROCESS SLIDE
            resiliency_analysis_slide_layout = prs.slide_layouts[slide_deck[8]] #CYBER RESILIENCY ANALYSIS SLIDE

            #CLASSIFICATION MARKINGS
            long_gen_marking = 'UNCLASSIFIED//FOR OFFICIAL USE ONLY'
            short_gen_marking = '(CUI)'
            gen_color = [45, 145, 45] #green
            if event.classification == 'UNCLASSIFIED':
                color = [45, 145, 45] #green
                if event.designation == 'CUI':
                    long_marking = 'CUI'
                    short_marking = '(CUI)'
                    if event.draft == 'YES':
                        slide_marking = 'CUI//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'CUI'
                elif event.designation == 'FOUO':
                    long_marking = 'UNCLASSIFIED//FOUO'
                    short_marking = '(U//FOUO)'
                    short_gen_marking = '(U//FOUO)'
                    if event.draft == 'YES':
                        slide_marking = 'UNCLASSIFIED//FOR OFFICIAL USE ONLY//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'UNCLASSIFIED//FOR OFFICIAL USE ONLY'
                else:
                    long_marking = 'UNCLASSIFIED'
                    short_marking = '(U)'
                    long_gen_marking = 'UNCLASSIFIED'
                    short_gen_marking = '(U)'
                    if event.draft == 'YES':
                        slide_marking = 'UNCLASSIFIED//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'UNCLASSIFIED'
            if event.classification == 'SECRET':
                color = [255, 0, 0] #red
                if event.designation == 'NOFORN':
                    long_marking = 'SECRET//NOFORN'
                    short_marking = '(S//NOFORN)'
                    if event.draft == 'YES':
                        slide_marking = 'SECRET//NOFORN//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'SECRET//NOFORN'
                else:
                    long_marking = 'SECRET'
                    short_marking = '(S)'
                    if event.draft == 'YES':
                        slide_marking = 'SECRET//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'SECRET'
            if event.classification == 'TOP SECRET':
                if event.designation == 'SCI':
                    color = [186, 186, 0] #yellow
                    long_marking = 'TOP SECRET//SCI'
                    short_marking = '(TS//SCI)'
                    if event.draft == 'YES':
                        slide_marking = 'TOP SECRET//SCI//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'TOP SECRET//SCI'
                else:
                    color = [255, 140, 0] #orange
                    long_marking = 'TOP SECRET'
                    short_marking = '(TS)'
                    if event.draft == 'YES':
                        slide_marking = 'TOP SECRET//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'TOP SECRET'

            def set_classification(body, marking, colors, alignment):
                tf = body.text_frame
                p = tf.paragraphs[0]
                if alignment == 'left':
                    p.alignment = PP_ALIGN.LEFT
                elif alignment == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = marking
                font = run.font
                font.name = 'Arial'
                font.size = Pt(8)
                font.bold = True
                font.color.rgb = RGBColor(colors[0], colors[1], colors[2])

            def set_distribution(body, text1, text2):
                tf = body.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.JUSTIFY
                run = p.add_run()
                run.text = text1
                font = run.font
                font.name = 'Arial'
                font.size = Pt(8)
                font.bold = True
                run = p.add_run()
                run.text = text2
                font = run.font
                font.name = 'Arial'
                font.size = Pt(8)
                font.bold = False
            
            def set_cui(body, text):
                tf = body.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = text
                font = run.font
                font.name = 'Arial'
                font.size = Pt(8)
                font.bold = False

            def set_title(shape, text):
                title = shape.title
                title.text = text

            def set_text(body, text):
                body.text = text
            
            def set_bullet(body, text, level):
                tf = body.text_frame
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                
            def set_findings(body, attrib, text, mode):
                tf = body.text_frame
                p = tf.add_paragraph()
                if mode == 2:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = attrib
                font = run.font
                font.bold = True
                run = p.add_run()
                run.text = str(text)
                font = run.font
                font.bold = False
                if mode == 0:
                    p = tf.add_paragraph()
            
            def set_info(body, text, f_bold, f_size, c_color):
                tf = body.text_frame
                p = tf.add_paragraph()
                p.text = text
                p.font.bold = f_bold
                p.font.size = Pt(f_size)
                p.font.color.rgb = RGBColor(c_color[0], c_color[1], c_color[2])
                p.alignment = PP_ALIGN.CENTER

            def _add_image(slide, placeholder_id, image_url):
                placeholder = slide.placeholders[placeholder_id]
                im = Image.open(image_url)
                #REMOVE DEFAULT IMAGE PLACEHOLDER
                image = slide.shapes[1]
                sp = image.element
                sp.getparent().remove(sp)
                #INSERT IMAGE WITH NO DEFAULT PLACEHOLDER
                im_width, im_height = im.size
                image_ratio = (im_width)/float(1.0) / (im_height)/float(1.0)
                fixed_width = 9.70
                fixed_height = 5.20
                fixed_ratio = (fixed_width)/float(1.0) / (fixed_height)/float(1.0)
                top = Inches(1.62)
                if image_ratio > fixed_ratio:
                    new_width = Inches(fixed_width)
                    new_height = Inches(fixed_width / image_ratio)
                    #ADJUST MARKINGS PLACEHOLDERS VERTICALLY
                    mid_height = (fixed_height - (fixed_width / image_ratio))/float(2.0)
                    #TOP
                    shapes.placeholders[17].top = Inches(1.24 + mid_height)
                    shapes.placeholders[17].left = Inches(2.44)
                    shapes.placeholders[17].width = Inches(2.5)
                    shapes.placeholders[17].height = Inches(0.39)
                    #BOTTOM
                    shapes.placeholders[18].top = Inches(6.81 - mid_height)
                    shapes.placeholders[18].left = Inches(2.44)
                    shapes.placeholders[18].width = Inches(2.5)
                    shapes.placeholders[18].height = Inches(0.39)
                    #ADJUST IMAGE VERTICAL POSITION
                    top = Inches(1.62 + mid_height)
                else:
                    new_width = Inches(fixed_height * image_ratio)
                    new_height = Inches(fixed_height)
                left = Inches(5) - (new_width/2)
                pic = shapes.add_picture(image_url, left, top, width=new_width, height=new_height)
                #ADJUST MARKINGS PLACEHOLDERS HORIZONTALLY
                #TOP
                shapes.placeholders[17].top = int(top - Inches(0.39))
                shapes.placeholders[17].left = int(left - Inches(0.12))
                shapes.placeholders[17].width = Inches(2.5)
                shapes.placeholders[17].height = Inches(0.39)
                #BOTTOM
                shapes.placeholders[18].top = int(top + new_height)
                shapes.placeholders[18].left = int(left + new_width - Inches(2.5) + Inches(0.12))
                shapes.placeholders[18].width = Inches(2.5)
                shapes.placeholders[18].height = Inches(0.39)
                #ADD FRAME TO SCREENSHOT
                line = pic.line
                line.color.rgb = RGBColor(0xFF, 0xDA, 0x3D)
                line.width = Pt(3)


            #COVER SLIDE
            if event.designation == 'CUI' or event.classification == 'SECRET' or event.classification == 'TOP SECRET':
                slide = prs.slides.add_slide(cover_slide_cui_layout)
            else:
                slide = prs.slides.add_slide(cover_slide_layout)
            shapes = slide.shapes
            #TITLE OF PROJECT
            if event.event_name != '':
                text = '(U) ' + event.event_name
            else:
                text = '(U) SUBTITLE GOES HERE'
            set_text(shapes.placeholders[12], text)
            if event.event_type == 'CVPA':
                set_text(shapes.placeholders[22], '(U) Emerging Results Brief (ERB)')
            else:
                set_text(shapes.placeholders[22], '(U) Cyber Resiliency Outbrief – Prevent, Mitigate, and Recover (PMR)')
            #NAME OF LEAD
            if event.lead_name != '':
                text = event.lead_name
            else:
                text = 'Name of Presenter'
            set_text(shapes.placeholders[13], text)
            #TITLE/RANK OF LEAD
            if event.lead_title != '':
                text = event.lead_title
            else:
                text = 'Rank/Title of Presenter'
            set_text(shapes.placeholders[14], text)
            #ORGANIZATION OF LEAD
            if event.lead_org != '':
                text = event.lead_org
            else:
                text = 'Organization of Presenter'
            set_text(shapes.placeholders[15], text)
            #DATE
            end_date = self.endDate.date()
            end_date = end_date.toString("dd MMM yyyy")
            if end_date != '':
                text = end_date
            else:
                text = 'DD MMM YYYY'
            set_text(shapes.placeholders[17], text)
            #CUI STATEMENT
            if event.designation == 'CUI' or event.classification == 'SECRET' or event.classification == 'TOP SECRET':
                office_symbol = event.office_symbol[0:10]
                text = 'Controlled by: ' + office_symbol + '\nCUI Category: DCRIT, Export Control\nDistribution/Dissemination Controls: D\nPOC: ' + event.lead_name + ', ' + event.lead_office
                set_cui(shapes.placeholders[24], text)
            #DISTRIBUTION STATEMENT
            dist_date = self.endDate.date().longMonthName(self.endDate.date().month()) + ' ' + str(self.endDate.date().year())
            text1 = '            DISTRIBUTION STATEMENT D. '
            text2 = 'Distribution authorized to the Department of Defense and U.S. DOD contractors only; administrative or operational use; export control; vulnerability information (' + dist_date + '). Other requests for this document shall be referred to Director, U.S. Army DEVCOM Analysis Center, White Sands Missile Range, NM 88002.'
            set_distribution(shapes.placeholders[23], text1, text2)
            #CLASSIFICATION - SLIDE
            set_classification(shapes.placeholders[20], slide_marking, color, 'center')
            set_classification(shapes.placeholders[21], slide_marking, color, 'center')


            if event.event_type == 'CVPA':
                #SCOPE SLIDE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) SCOPE')
                #BODY WITH BULLETS
                if scope_list:
                    prev_bul = 0
                    for i in range(len(scope_list)):
                        if '* ' in scope_list[i]:
                            new_string = short_marking + ' ' + scope_list[i].replace('* ', '')
                            if prev_bul == 1:
                                set_bullet(shapes.placeholders[1], '', 0)
                                set_bullet(shapes.placeholders[1], new_string, 0)
                            else:
                                set_text(shapes.placeholders[1], new_string)
                            prev_bul = 1
                        if '    -' in scope_list[i]:
                            new_string = short_marking + ' ' + scope_list[i].replace('    -', '')
                            set_bullet(shapes.placeholders[1], new_string, 1)
                else:
                    set_text(shapes.placeholders[1], short_marking + ' <<Add Scope description as needed>>')
                    set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                    set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                    set_bullet(shapes.placeholders[1], '', 0)
                    set_bullet(shapes.placeholders[1], short_marking + ' ', 0)
                    set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                    set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #SYSTEM UNDER TEST
                num_diagrams = len(diagrams_list)
                current_diagram = 0
                while(num_diagrams >= 0):
                    slide = prs.slides.add_slide(picture_slide_layout)
                    shapes = slide.shapes
                    #TITLE
                    set_title(shapes, '(U) SYSTEM UNDER TEST')
                    #INSERT FIGURE
                    if num_diagrams != 0:
                        i_file = cwd + '/' + diagrams_list[current_diagram]
                        if os.path.isfile(i_file):
                            _add_image(slide, 10, i_file)
                    if num_diagrams == 0:
                        #ADJUST MARKINGS PLACEHOLDERS HORIZONTALLY
                        #TOP
                        shapes.placeholders[17].top = int(Inches(1.64 - 0.39))
                        shapes.placeholders[17].left = int(Inches(0.88 - 0.12))
                        shapes.placeholders[17].width = Inches(2.5)
                        shapes.placeholders[17].height = Inches(0.39)
                        #BOTTOM
                        shapes.placeholders[18].top = int(Inches(1.64 + 5.14))
                        shapes.placeholders[18].left = int(Inches(0.88 + 8.27 - 2.5 + 0.12))
                        shapes.placeholders[18].width = Inches(2.5)
                        shapes.placeholders[18].height = Inches(0.39)
                    #CLASSIFICATION - FIGURE
                    tmp_color = color
                    if diagram_marking == 'UNCLASSIFIED':
                        d_marking = 'UNCLASSIFIED'
                        color = [45, 145, 45] #green
                    else:
                        d_marking = long_marking
                    set_classification(shapes.placeholders[17], d_marking, color, 'left')
                    set_classification(shapes.placeholders[18], d_marking, color, 'right')
                    color = tmp_color
                    #CLASSIFICATION - SLIDE
                    set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                    set_classification(shapes.placeholders[21], slide_marking, color, 'center')
                    num_diagrams -= 1
                    current_diagram += 1
                    if num_diagrams == 0:
                        break


                #AGENDA - EXECUTED ACTIVITIES SLIDE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) AGENDA - EXECUTED ACTIVITIES')
                #WORK DAYS
                work_days = []
                day_1 = self.startDate.date()
                str_day_1 = day_1.toString("dd MMM yyyy")
                work_days.append(str_day_1)
                for i in range(1,5):
                    next_day = day_1.addDays(i)
                    str_next_day = next_day.toString("dd MMM yyyy")
                    work_days.append(str_next_day)
                #DAY 1
                set_text(shapes.placeholders[1], short_gen_marking + ' ' + work_days[0])
                text = '(U) In-processing, setup, network connectivity testing, discovery and enumeration scans, started collecting DOT&E metrics, and began the penetration test.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                #DAY 2
                set_bullet(shapes.placeholders[1], short_gen_marking + ' ' + work_days[1], 0)
                text = '(U) Continuation of the penetration test and DOT&E metric collection.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                #DAY 3
                set_bullet(shapes.placeholders[1], short_gen_marking + ' ' + work_days[2], 0)
                text = '(U) Continuation of the penetration test, DOT&E metric collection, and Personnel Interviews.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                #DAY 4
                set_bullet(shapes.placeholders[1], short_gen_marking + ' ' + work_days[3], 0)
                text = '(U) Continuation of the penetration test and DOT&E metric collection.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                #DAY 5
                set_bullet(shapes.placeholders[1], short_gen_marking + ' ' + work_days[4], 0)
                text = '(U) Completed the penetration test, performed the system cleanup and restoration, data consolidation, and backup. Performed the Emerging Results Brief (ERB) presentation to stakeholders.'
                set_bullet(shapes.placeholders[1], text, 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #PENETRATION TESTING ACTIVITIES
                slide = prs.slides.add_slide(pen_process_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) ACTIVITIES COMPLETION STATUS')
                #CLASSIFICATION - FIGURE
                set_classification(shapes.placeholders[17], 'CUI', gen_color, 'left')
                set_classification(shapes.placeholders[22], 'CUI', gen_color, 'right')
                #ADJUST MARKINGS PLACEHOLDERS HORIZONTALLY
                #TOP
                shapes.placeholders[17].top = Inches(0.90)
                shapes.placeholders[17].left = Inches(0.22)
                shapes.placeholders[17].width = Inches(2.5)
                shapes.placeholders[17].height = Inches(0.39)
                #BOTTOM
                shapes.placeholders[22].top = Inches(6.94)
                shapes.placeholders[22].left = Inches(7.20)
                shapes.placeholders[22].width = Inches(2.5)
                shapes.placeholders[22].height = Inches(0.39)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #PENETRATION TESTING PROCESS SLIDE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) PENETRATION TESTING PROCESS')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '(U) Characterization of key cyber terrain and attack vector generation')
                text = '(U) Documentation review, OSINT, site visit, staff interview, identify cyber postures, and develop attack vectors.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Discovery and Enumeration Scans'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) Map network, automated scanning for well-known weaknesses.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Penetration Testing'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) Manual probing, exploration, data pillaging, lateral movement.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Risk Analysis'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) Assess impact to confidentiality, integrity, and availability.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Mitigation and Risk Reduction Strategies'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) Develop and provide potential mitigation and risk reduction strategies to the discovered vulnerabilities.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Follow-on Testing'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) After mitigations are implemented, re-test to ensure the fixes are effective and do not introduce new vulnerabilities.'
                set_bullet(shapes.placeholders[1], text, 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #POSTURES SLIDE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) POSTURES')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '(U) Findings in this ERB constitute raw results and the technical risk analysis has not been determined.')
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) All technical findings assume some level of physical or logical access to the assets.'
                set_bullet(shapes.placeholders[1], text, 0)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Each finding will be from a specific posture. We define this postures to be as follow:'
                set_bullet(shapes.placeholders[1], text, 0)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Insider – is a person with legitimate access to the system, both logical (credentialed user) and physical or remote access.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Nearsider – physical access is provided to the target network and system, but with no credentials given.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Outsider – is a person without legitimate physical and logical access to the system under test and it is placed outside the accreditation boundary. The outsider posture is normally portrayed by an actor pivoting off a system that is legitimate connected external vectors such as SIPRNet, or Sensors.'
                set_bullet(shapes.placeholders[1], text, 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #FINDINGS TABLE SLIDE - 12 FINDINGS PER SLIDE MAX
                limit = 12
                per_slide = []
                if len(order_list) > 0:
                    reminder = len(order_list)%limit
                    complete = int(len(order_list)/limit)
                    if complete > 0:
                        for i in range(complete):
                            per_slide.append(12)
                    if reminder != 0:
                        per_slide.append(reminder)
                    for i in range(len(per_slide)):
                        slide = prs.slides.add_slide(picture_slide_layout)
                        shapes = slide.shapes
                        #REMOVE DEFAULT IMAGE PLACEHOLDER
                        image = slide.shapes[1]
                        sp = image.element
                        sp.getparent().remove(sp)
                        #CLASSIFICATION - FIGURE
                        set_classification(shapes.placeholders[17], long_marking, color, 'left')
                        set_classification(shapes.placeholders[18], long_marking, color, 'right')
                        #TITLE
                        set_title(shapes, '(U) Table of Findings')
                        offset = 2.7 - (per_slide[i]*0.1)
                        t_left = Inches(1.2)
                        t_offset = offset - 0.28
                        t_top = Inches(t_offset)
                        #TABLE
                        rows = per_slide[i]+1
                        cols = 2
                        left = Inches(1.2)
                        top = Inches(offset)
                        width = Inches(6.0)
                        height = Inches(0.8)
                        table_1 = shapes.add_table(rows, cols, left, top, width, height).table
                        #SET COLUMN WIDTHS
                        table_1.columns[0].width = Inches(0.5)
                        table_1.columns[1].width = Inches(7.0)
                        #SET COLUMN HEADINGS
                        table_1.cell(0,1).text = 'Findings'
                        #WRITE BODY CELLS
                        for j in range(per_slide[i]):
                            index = j + (i*limit)
                            table_1.cell(j+1,0).text = str(index+1)
                            table_1.cell(j+1,1).text = findings[order_list[index]].title
                            if len(table_1.cell(j+1,1).text) > 62:
                                cell = table_1.rows[j+1].cells[1]
                                paragraph = cell.text_frame.paragraphs[0]
                                paragraph.font.size = Pt(10)
                        #ADJUST MARKINGS PLACEHOLDERS
                        #TOP
                        shapes.placeholders[17].top = int(t_top - Inches(0.05))
                        shapes.placeholders[17].left = int(t_left - Inches(0.07))
                        shapes.placeholders[17].width = Inches(2.5)
                        shapes.placeholders[17].height = Inches(0.39)
                        #BOTTOM
                        t_top = offset + 0.4*(per_slide[i]+1)
                        shapes.placeholders[18].top = Inches(t_top - 0.05)
                        shapes.placeholders[18].left = int(t_left + Inches(7.5/2) + Inches(2.5/2) + Inches(0.09))
                        shapes.placeholders[18].width = Inches(2.5)
                        shapes.placeholders[18].height = Inches(0.39)
                        #CLASSIFICATION
                        set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                        set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #FINDINGS - DESCRIPTION & SCREENSHOTS
                for i in order_list:
                    #FINDING SLIDE
                    slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                    shapes = slide.shapes
                    #TITLE
                    set_title(shapes, short_marking + ' ' + findings[i].title)
                    #POSTURE
                    set_findings(shapes.placeholders[1], short_marking + ' Posture: ', findings[i].posture.capitalize(), 0)
                    #AFFECTED SYSTEMS
                    set_findings(shapes.placeholders[1], short_marking + ' Affected System(s): ', findings[i].hosts, 0)
                    #ISSUES
                    set_findings(shapes.placeholders[1], short_marking + ' Issue(s): ', findings[i].issues, 1)
                    #MITIGATION
                    if findings[i].include_mitigation == 'yes':
                        set_findings(shapes.placeholders[1], short_marking + ' Mitigation: ', findings[i].mitigation, 2)
                    #CLASSIFICATION
                    set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                    set_classification(shapes.placeholders[21], slide_marking, color, 'center')
                    #SCREENSHOTS
                    if findings[i].screenshots != None:
                        for j in findings[i].screenshots:
                            ss_folder = findings[i].folder
                            i_file = ss_folder + j
                            slide = prs.slides.add_slide(picture_slide_layout)
                            shapes = slide.shapes
                            #TITLE
                            set_title(shapes, short_marking + ' ' + findings[i].title)
                            #FIGURE
                            _add_image(slide, 10, i_file)
                            #CLASSIFICATION - FIGURE
                            set_classification(shapes.placeholders[17], long_marking, color, 'left')
                            set_classification(shapes.placeholders[18], long_marking, color, 'right')
                            #CLASSIFICATION - SLIDE
                            set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                            set_classification(shapes.placeholders[21], slide_marking, color, 'center')


            if event.event_type == 'PMR':
                #CYBER RESILIENCY OBJECTIVE SLIDE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) CYBER RESILIENCY OBJECTIVE')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '(CUI) In support of <<SYSTEM / PROGRAM>> Evaluation and Assessments')
                text = '(CUI) Data collected during the event will inform about the cyber resilience posture.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Guided by Director, Operational Test and Evaluation (DOT&E) Memo'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) “Procedure For Operational Test And Evaluation Of Cybersecurity In Acquisition Programs” April 3, 2018'
                set_bullet(shapes.placeholders[1], text, 1)
                text = '(U) Prevent – The ability to protect critical mission functions from cyber threats.'
                set_bullet(shapes.placeholders[1], text, 2)
                text = '(U) Mitigate – The ability to detect and respond to cyber-attacks, and assess resilience to survive attacks and complete critical missions and tasks.'
                set_bullet(shapes.placeholders[1], text, 2)
                text = '(U) Recover – The resilience to recover from cyber-attacks and prepare mission systems”'
                set_bullet(shapes.placeholders[1], text, 2)
                text1 = '(CUI) OUR GOAL:'
                text2 = ' Provide the data and analysis on the cyber resiliency of the system to stakeholders and determine how the overall Defense Cyber Operations Team (DCOT) performed in a cyber-contested environment. Assist stakeholders to understand and have a clear picture of gaps and their strengths in DCO. Assist the program by proposing solutions and mitigations to further increase the cyber-robustness of the systems.'
                tf = shapes.placeholders[1].text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = text1
                font = run.font
                font.size = Pt(14)
                font.bold = True
                run = p.add_run()
                run.text = text2
                font = run.font
                font.size = Pt(14)
                font.bold = False
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #PMR CRITERIA - PREVENT
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) PMR CRITERIA')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '(U) PREVENT:')
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '(U) Defined as an action that is typically proactive and implemented before the occurrence of a cyber-threat activity.', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '(U) An activity that is part of an attack vector executed by TSMO that does not meet the intended objective/effect due to a preventative mechanism in place will result in a successful prevention.', 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #PMR CRITERIA - MITIGATE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) PMR CRITERIA')
                #BODY WITH BULLETS
                text1 = '(U) MITIGATE –'
                text2 = ' Activities are divided into two subcategories:'
                tf = shapes.placeholders[1].text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = text1
                font = run.font
                font.size = Pt(16)
                font.bold = True
                run = p.add_run()
                run.text = text2
                font = run.font
                font.size = Pt(16)
                font.bold = False
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '(U) DETECT', 1)
                set_bullet(shapes.placeholders[1], '(U) A successful detect is accomplished when a defender/operator acknowledges the technical/system detect activity and consequently creates a report (cyber incident report/helpdesk ticket), this can be triggered by an automated alert or manual review of logs. A human detect can also be accomplished by an operator/defender detecting abnormal activity on their system without the need for a network defense or monitoring tool.', 2)
                set_bullet(shapes.placeholders[1], '(U) REACT', 1)
                set_bullet(shapes.placeholders[1], '(U) A successful reaction is defined as force or action that successfully counteracts, hinders, thwarts, and/or mitigates the cyber threat action (i.e. denying the threat access to a service, port, and/or host).', 2)
                set_bullet(shapes.placeholders[1], '', 0)
                text1 = '(U) Note:'
                text2 = ' Mitigate metrics as specified in Attachment C of the DOT&E document will be based off of information extracted from logs, incident reports, observer logs, help desk tickets, and other products from the defender that will be parsed for the information requested under the measurements column of the mitigate actions. PMR observer logs will be used to fill-in any gaps left by the defender when necessary.'
                tf = shapes.placeholders[1].text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = text1
                font = run.font
                font.size = Pt(14)
                font.bold = True
                run = p.add_run()
                run.text = text2
                font = run.font
                font.size = Pt(14)
                font.bold = False
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #PMR CRITERIA - RECOVER
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) PMR CRITERIA')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '(U) RECOVER:')
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '(U) Recover activities are those taken by operators and/or network defenders to restore mission and/or technical capabilities to continue operations after a degradation of such capabilities.', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '(U) Recover activities include, but are not limited to, re-imaging a host and/or failover to alternate sites.', 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #CYBER RESILIENCY ANALYSIS
                slide = prs.slides.add_slide(resiliency_analysis_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) CYBER RESILIENCY ANALYSIS')
                #CLASSIFICATION - FIGURE
                set_classification(shapes.placeholders[17], 'CUI', gen_color, 'left')
                set_classification(shapes.placeholders[22], 'CUI', gen_color, 'right')
                #ADJUST MARKINGS PLACEHOLDERS HORIZONTALLY
                #TOP
                shapes.placeholders[17].top = Inches(1.80)
                shapes.placeholders[17].left = Inches(0.22)
                shapes.placeholders[17].width = Inches(2.5)
                shapes.placeholders[17].height = Inches(0.39)
                #BOTTOM
                shapes.placeholders[22].top = Inches(5.95)
                shapes.placeholders[22].left = Inches(7.20)
                shapes.placeholders[22].width = Inches(2.5)
                shapes.placeholders[22].height = Inches(0.39)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')

                #WHITE CARDS
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) WHITECARDS')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], short_marking + ' The following Whitecards were granted to the TSMO AA Team:')
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' (1) mm/dd/yyyy @ hh:mm – <<Add Whitecard description as needed>>', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' (2) mm/dd/yyyy @ hh:mm – <<Add Whitecard description as needed>>', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' (3) mm/dd/yyyy @ hh:mm – <<Add Whitecard description as needed>>', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' (4) mm/dd/yyyy @ hh:mm – <<Add Whitecard description as needed>>', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' (5) mm/dd/yyyy @ hh:mm – <<Add Whitecard description as needed>>', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')

            #OVERALL OBSERVATIONS
            slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
            shapes = slide.shapes
            #TITLE
            if event.event_type == 'CVPA':
                set_title(shapes, '(U) OVERALL OBSERVATIONS')
            else:
                set_title(shapes, '(U) CYBER RESILIENCY TEAM OBSERVATIONS')
            #BODY WITH BULLETS
            set_text(shapes.placeholders[1], short_marking + ' System Strengths')
            if strengths_list:
                for i in range(len(strengths_list)):
                    if '* ' in strengths_list[i]:
                        new_string = short_marking + ' ' + strengths_list[i].replace('* ', '')
                        set_bullet(shapes.placeholders[1], new_string, 1)
            else:
                set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
            set_bullet(shapes.placeholders[1], '', 0)
            set_bullet(shapes.placeholders[1], short_marking + ' System Weaknesses', 0)
            if weaknesses_list:
                for i in range(len(weaknesses_list)):
                    if '* ' in weaknesses_list[i]:
                        new_string = short_marking + ' ' + weaknesses_list[i].replace('* ', '')
                        set_bullet(shapes.placeholders[1], new_string, 1)
            else:
                set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
            if event.event_type == 'CVPA':
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' Overall Mitigations', 0)
                if mitigations_list:
                    for i in range(len(mitigations_list)):
                        if '* ' in mitigations_list[i]:
                            new_string = short_marking + ' ' + mitigations_list[i].replace('* ', '')
                            set_bullet(shapes.placeholders[1], new_string, 1)
                else:
                    set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                    set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
            #CLASSIFICATION - SLIDE
            set_classification(shapes.placeholders[20], slide_marking, color, 'center')
            set_classification(shapes.placeholders[21], slide_marking, color, 'center')


            if event.event_type == 'CVPA':
                #POST ASSESSMENT REPORTING
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) POST ASSESSMENT REPORTING')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '')
                set_bullet(shapes.placeholders[1], '', 0)
                if deliverables:
                    if deliverables[0] == 'YES': #ERB
                        set_bullet(shapes.placeholders[1], short_gen_marking + ' Emerging Results Brief (ERB)', 0)
                        if prefix == 'VOF_':
                            set_bullet(shapes.placeholders[1], '(U) Findings Status --> Open / Closed', 1)
                        else:
                            set_bullet(shapes.placeholders[1], '(U) List of findings with minimal analysis', 1)
                            set_bullet(shapes.placeholders[1], '(U) Overall assessment objective completion status', 1)
                        set_bullet(shapes.placeholders[1], '', 0)
                        set_bullet(shapes.placeholders[1], '', 0)
                    if deliverables[1] == 'YES': #ARM
                        set_bullet(shapes.placeholders[1], short_gen_marking + ' Assessment Results Matrix (ARM)', 0)
                        set_bullet(shapes.placeholders[1], '(U) List of findings with technical risk levels', 1)
                        set_bullet(shapes.placeholders[1], '', 0)
                        set_bullet(shapes.placeholders[1], '', 0)
                    if deliverables[3] == 'MEMO': #MEMO OR REPORT
                        new_string = 'Memorandum'
                    else:
                        new_string = 'Report'
                    if deliverables[2] == 'YES': #MEMO/REPORT --> # of days
                        set_bullet(shapes.placeholders[1], short_gen_marking + ' Technical ' + new_string, 0)
                        set_bullet(shapes.placeholders[1], '(U) ~' + deliverables[4] + ' working days', 1)
                else:
                    set_bullet(shapes.placeholders[1], short_gen_marking + ' Emerging Results Brief (ERB)', 0)
                    set_bullet(shapes.placeholders[1], '(U) List of findings with minimal analysis', 1)
                    set_bullet(shapes.placeholders[1], '(U) Overall assessment objective completion status', 1)
                    set_bullet(shapes.placeholders[1], '', 0)
                    set_bullet(shapes.placeholders[1], '', 0)
                    set_bullet(shapes.placeholders[1], short_gen_marking + ' Assessment Results Matrix (ARM)', 0)
                    set_bullet(shapes.placeholders[1], '(U) List of findings with technical risk levels', 1)
                    set_bullet(shapes.placeholders[1], '', 0)
                    set_bullet(shapes.placeholders[1], '', 0)
                    set_bullet(shapes.placeholders[1], short_gen_marking + ' Technical Report', 0)
                    set_bullet(shapes.placeholders[1], '(U) Option A --> Technical Memorandum ~30 working days', 1)
                    set_bullet(shapes.placeholders[1], '(U) Option B --> Published Report ~90 working days', 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


            #LAST SLIDE - CONTACT INFORMATION SLIDE
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            #TITLE
            set_title(shapes, '(U) Contact Information')
            left = Inches(1)
            top = Inches(2)
            width = height = Inches(8)
            #BODY WITH NO BULLETS
            txBox = slide.shapes.add_textbox(left, top, width, height)
            set_text(txBox.text_frame, '')
            if event.lead_name != '':
                set_info(txBox, event.lead_name, True, 28, c_color)
            else:
                set_info(txBox, '<<TEAM LEAD>>', True, 28, c_color)
            if event.lead_nipr != '':
                n_email = 'UNCLASSIFIED: ' + event.lead_nipr
            else:
                n_email = 'UNCLASSIFIED: xxxxxxxx@army.mil'
            set_info(txBox, n_email, False, 20, c_color)
            if event.lead_sipr != '':
                s_email = 'SIPR: ' + event.lead_sipr
            else:
                s_email = 'SIPR: xxxxxxxx@mail.smil.mil'
            set_info(txBox, s_email, False, 20, c_color)
            if event.lead_office != '':
                o_number = 'O: ' + event.lead_office
                set_info(txBox, o_number, False, 20, c_color)
            if event.lead_mobile != '':
                m_number = 'M: ' + event.lead_mobile
                set_info(txBox, m_number, False, 20, c_color)
            #CLASSIFICATION - SLIDE
            set_classification(shapes.placeholders[20], slide_marking, color, 'center')
            set_classification(shapes.placeholders[21], slide_marking, color, 'center')


            #SAVE DOCUMENT
            o_file = desktop_dir + '/' + pptx_file
            prs.save(o_file)
            #MESSAGE BOX!
            msg = QMessageBox()
            msg.setIcon(QMessageBox.Information)
            msg.setWindowTitle("PowerPoint")
            txt_msg = "PowerPoint successfully created at " + o_file + "\nThank you... Good bye!"
            msg.setText(txt_msg)
            x = msg.exec_()
            #EXIT
            app.quit()


        #CHECK IF SUT DIAGRAMS EXIST:
        def purge_diagrams():
            global diagrams_delete
            for i in diagrams_delete:
                if os.path.isfile(i):
                    cmd = 'rm ' + cwd + '/' + i
                    os.system(cmd)


        #SYSTEM UNDER TEST (SUT) PAGE
        def write_sut():
            global scope_list, diagram_marking, diagrams_list, strengths_list, weaknesses_list, mitigations_list, deliverables
            #CREATE NEW XML FILE
            xml_file = cwd + '/sut_info.xml'
            def indent(elem, level=0):
                i = "\n" + level*"    "
                if len(elem):
                    if not elem.text or not elem.text.strip():
                        elem.text = i + "    "
                    if not elem.tail or not elem.tail.strip():
                        elem.tail = i
                    for elem in elem:
                        indent(elem, level+1)
                    if not elem.tail or not elem.tail.strip():
                        elem.tail = i
                else:
                    if level and (not elem.tail or not elem.tail.strip()):
                        elem.tail = i
            #CREATE FILE STRUCTURE
            root = ET.Element('data')
            #CREATE SUT
            xml_sut = ET.SubElement(root, 'sut')
            xml_sut.set('uid', str(0))
            xml_scope_list = ET.SubElement(xml_sut, 'scope')
            xml_scope_list.text = str(scope_list)
            xml_diagram_marking = ET.SubElement(xml_sut, 'diagram')
            xml_diagram_marking.text = str(diagram_marking)
            xml_diagrams_list = ET.SubElement(xml_sut, 'diagrams_list')
            xml_diagrams_list.text = str(diagrams_list)
            purge_diagrams()
            xml_strengths_list = ET.SubElement(xml_sut, 'strengths')
            xml_strengths_list.text = str(strengths_list)
            xml_weaknesses_list = ET.SubElement(xml_sut, 'weaknesses')
            xml_weaknesses_list.text = str(weaknesses_list)
            xml_mitigations_list = ET.SubElement(xml_sut, 'mitigations')
            xml_mitigations_list.text = str(mitigations_list)
            xml_deliverables = ET.SubElement(xml_sut, 'deliverables')
            xml_deliverables.text = str(deliverables)
            #WRITING XML
            indent(root)
            tree = ET.ElementTree(root)
            tree.write(xml_file, encoding='utf-8', xml_declaration=True)


        #ADD SCOPE TO LIST
        def addScope():
            s_input = self.inputScope.text()
            if s_input != "":
                current_list = [str(self.listScope.item(i).text()) for i in range(self.listScope.count())]
                if len(current_list) == 0:
                    self.bulletButton.setChecked(True)
                if self.bulletButton.isChecked() == True:
                    mod_input = "* " + s_input
                if self.subbulletButton.isChecked() == True:
                    mod_input = "    - " + s_input
                self.listScope.addItem(mod_input)
                self.inputScope.clear()


        #DELETE SCOPE FROM LIST
        def delScope():
            if self.listScope.selectedItems() != []:
                currentRow = self.listScope.currentRow()
                currentItem = self.listScope.item(currentRow).text()
                if '* ' in currentItem:
                    item = self.listScope.takeItem(self.listScope.currentRow())
                    current_list = [str(self.listScope.item(i).text()) for i in range(self.listScope.count())]
                    offset = 0
                    while(True):
                        if currentRow+offset < len(current_list):
                            if '    - ' in current_list[currentRow+offset]:
                                offset+=1
                            else:
                                break
                        else:
                            break
                    for i in range(offset):
                        item = self.listScope.takeItem(self.listScope.currentRow())
                if '    - ' in currentItem:
                    item = self.listScope.takeItem(self.listScope.currentRow())


        #ADD DIAGRAM TO LIST
        def addDiagram():
            global f_folder, diagrams_list
            if f_folder == '':
                f_folder = desktop_dir
            fname, _filter = QtWidgets.QFileDialog.getOpenFileName(None, "Import Image", f_folder, "*.png *.jpg *.bmp *.gif *.jpeg *.tif *.tiff")
            if fname != '':
                current_list = [str(self.listSUT_diagrams.item(i).text()) for i in range(self.listSUT_diagrams.count())]
                aux = fname.rsplit('/', 1)
                f_folder = aux[0]
                d_name = aux[-1].split('.')
                d_name = d_name[0]
                msg = 'cp ' + str(fname) + ' ' + cwd + '/' + d_name
                os.system(msg)
                self.listSUT_diagrams.addItem(d_name)
                diagrams_list = [str(self.listSUT_diagrams.item(i).text()) for i in range(self.listSUT_diagrams.count())]
                self.listSUT_diagrams.item(len(current_list)).setSelected(True)
                self.listSUT_diagrams.setCurrentRow(len(current_list))
                self.sutPreview_image.setText("")
                self.sutPreview_image.setPixmap(QtGui.QPixmap(fname))
                self.sutPreview_image.setScaledContents(True)


        #DELETE DIAGRAM FROM LIST
        def delDiagram():
            global diagrams_list, diagrams_delete
            if self.listSUT_diagrams.selectedItems() != []:
                currentRow = self.listSUT_diagrams.currentRow()
                currentItem = self.listSUT_diagrams.item(currentRow).text()
                diagrams_delete.append(currentItem)
                item = self.listSUT_diagrams.takeItem(self.listSUT_diagrams.currentRow())
                diagrams_list = [str(self.listSUT_diagrams.item(i).text()) for i in range(self.listSUT_diagrams.count())]


        #MOVE ONE DIAGRAM UP THE LIST
        def diagramsUp():
            global diagrams_list
            if self.listSUT_diagrams.selectedItems() != []:
                currentRow = self.listSUT_diagrams.currentRow()
                currentItem = self.listSUT_diagrams.takeItem(currentRow)
                self.listSUT_diagrams.insertItem(currentRow - 1, currentItem)
                self.listSUT_diagrams.setCurrentItem(currentItem) 
                diagrams_list = [str(self.listSUT_diagrams.item(i).text()) for i in range(self.listSUT_diagrams.count())]
                preview_SUT()


        #MOVE ONE DIAGRAM DOWN THE LIST
        def diagramsDown():
            global diagrams_list
            if self.listSUT_diagrams.selectedItems() != []:
                currentRow = self.listSUT_diagrams.currentRow()
                currentItem = self.listSUT_diagrams.takeItem(currentRow)
                self.listSUT_diagrams.insertItem(currentRow + 1, currentItem)
                self.listSUT_diagrams.setCurrentItem(currentItem) 
                diagrams_list = [str(self.listSUT_diagrams.item(i).text()) for i in range(self.listSUT_diagrams.count())]
                preview_SUT()


        #DISPLAY IMAGES ON "SUT PREVIEW"
        def preview_SUT():
            if self.listSUT_diagrams.selectedItems() != []:
                item = self.listSUT_diagrams.selectedIndexes()[0]
                image_num = item.row()
                image_name = diagrams_list[image_num]
                image_file = cwd + '/'+ image_name
                self.sutPreview_image.setText("")
                self.sutPreview_image.setPixmap(QtGui.QPixmap(image_file))
                self.sutPreview_image.setScaledContents(True)
            else:
                self.sutPreview_image.setText("NO PREVIEW")


        #ADD STRENGTHS TO LIST
        def addStrengths():
            s_input = self.inputStrength.text()
            if s_input != "":
                mod_input = "* " + s_input
                self.listStrengths.addItem(mod_input)
                self.inputStrength.clear()


        #DELETE STRENGTHS FROM LIST
        def delStrengths():
            if self.listStrengths.selectedItems() != []:
                item = self.listStrengths.takeItem(self.listStrengths.currentRow())


        #MOVE ONE STRENGTH UP THE LIST
        def strengthsUp():
            if self.listStrengths.selectedItems() != []:
                currentRow = self.listStrengths.currentRow()
                currentItem = self.listStrengths.takeItem(currentRow)
                self.listStrengths.insertItem(currentRow - 1, currentItem)
                self.listStrengths.setCurrentItem(currentItem) 
                current_list = [str(self.listStrengths.item(i).text()) for i in range(self.listStrengths.count())]


        #MOVE ONE STRENGTH DOWN THE LIST
        def strengthsDown():
            if self.listStrengths.selectedItems() != []:
                currentRow = self.listStrengths.currentRow()
                currentItem = self.listStrengths.takeItem(currentRow)
                self.listStrengths.insertItem(currentRow + 1, currentItem)
                self.listStrengths.setCurrentItem(currentItem) 
                current_list = [str(self.listStrengths.item(i).text()) for i in range(self.listStrengths.count())]


        #ADD WEAKNESSES TO LIST
        def addWeaknesses():
            s_input = self.inputWeakness.text()
            if s_input != "":
                mod_input = "* " + s_input
                self.listWeaknesses.addItem(mod_input)
                self.inputWeakness.clear()


        #DELETE WEAKNESSES FROM LIST
        def delWeaknesses():
            if self.listWeaknesses.selectedItems() != []:
                item = self.listWeaknesses.takeItem(self.listWeaknesses.currentRow())


        #MOVE ONE WEAKNESS UP THE LIST
        def weaknessesUp():
            if self.listWeaknesses.selectedItems() != []:
                currentRow = self.listWeaknesses.currentRow()
                currentItem = self.listWeaknesses.takeItem(currentRow)
                self.listWeaknesses.insertItem(currentRow - 1, currentItem)
                self.listWeaknesses.setCurrentItem(currentItem) 
                current_list = [str(self.listWeaknesses.item(i).text()) for i in range(self.listWeaknesses.count())]


        #MOVE ONE WEAKNESS DOWN THE LIST
        def weaknessesDown():
            if self.listWeaknesses.selectedItems() != []:
                currentRow = self.listWeaknesses.currentRow()
                currentItem = self.listWeaknesses.takeItem(currentRow)
                self.listWeaknesses.insertItem(currentRow + 1, currentItem)
                self.listWeaknesses.setCurrentItem(currentItem) 
                current_list = [str(self.listWeaknesses.item(i).text()) for i in range(self.listWeaknesses.count())]


        #ADD MITIGATIONS TO LIST
        def addMitigations():
            s_input = self.inputMitigation.text()
            if s_input != "":
                mod_input = "* " + s_input
                self.listMitigations.addItem(mod_input)
                self.inputMitigation.clear()


        #DELETE MITIGATIONS FROM LIST
        def delMitigations():
            if self.listMitigations.selectedItems() != []:
                item = self.listMitigations.takeItem(self.listMitigations.currentRow())


        #MOVE ONE MITIGATION UP THE LIST
        def mitigationsUp():
            if self.listMitigations.selectedItems() != []:
                currentRow = self.listMitigations.currentRow()
                currentItem = self.listMitigations.takeItem(currentRow)
                self.listMitigations.insertItem(currentRow - 1, currentItem)
                self.listMitigations.setCurrentItem(currentItem) 
                current_list = [str(self.listMitigations.item(i).text()) for i in range(self.listMitigations.count())]


        #MOVE ONE MITIGATION DOWN THE LIST
        def mitigationsDown():
            if self.listMitigations.selectedItems() != []:
                currentRow = self.listMitigations.currentRow()
                currentItem = self.listMitigations.takeItem(currentRow)
                self.listMitigations.insertItem(currentRow + 1, currentItem)
                self.listMitigations.setCurrentItem(currentItem) 
                current_list = [str(self.listMitigations.item(i).text()) for i in range(self.listMitigations.count())]


        #GET SYSTEM DATA FROM GUI
        def system_data():
            global scope_list, diagram_marking, diagrams_list, strengths_list, weaknesses_list, mitigations_list, deliverables
            scope_list = [str(self.listScope.item(i).text()) for i in range(self.listScope.count())]
            if len(scope_list) == 1 and scope_list[0] == '':
                scope_list.clear()
            if self.unclassifiedButton.isChecked():
                diagram_marking = 'UNCLASSIFIED'
            if self.definedButton.isChecked():
                diagram_marking = self.definedButton.text()
            diagrams_list = [str(self.listSUT_diagrams.item(i).text()) for i in range(self.listSUT_diagrams.count())]
            strengths_list = [str(self.listStrengths.item(i).text()) for i in range(self.listStrengths.count())]
            if len(strengths_list) == 1 and strengths_list[0] == '':
                strengths_list.clear()
            weaknesses_list = [str(self.listWeaknesses.item(i).text()) for i in range(self.listWeaknesses.count())]
            if len(weaknesses_list) == 1 and weaknesses_list[0] == '':
                weaknesses_list.clear()
            mitigations_list = [str(self.listMitigations.item(i).text()) for i in range(self.listMitigations.count())]
            if len(mitigations_list) == 1 and mitigations_list[0] == '':
                mitigations_list.clear()
            deliverables = []
            if self.ERBcheckBox.isChecked() == True:
                deliverables.append('YES')
            else:
                deliverables.append('NO')
            if self.ARMcheckBox.isChecked() == True:
                deliverables.append('YES')
            else:
                deliverables.append('NO')
            if self.reportcheckBox.isChecked() == True:
                deliverables.append('YES')
            else:
                deliverables.append('NO')
            if self.memoButton.isChecked() == True:
                deliverables.append('MEMO')
            if self.reportButton.isChecked() == True:
                deliverables.append('REPORT')
            deliverables.append(self.inputDays.text())


        def sutData():
            system_data()
            write_sut()


        #SWITCH BETWEEN PAGES USING STACKED WIDGETS / CREATE PMR SLIDES
        def save_main():
            global event, folder_list
            if self.cvpaButton.isChecked() == True:
                event_type = 'CVPA'
            else:
                event_type = 'PMR'
            lead_name = self.leadName.text() #lead name
            lead_title = self.leadTitle.text() #lead title
            lead_org = self.leadOrg.text() #lead org
            lead_nipr = self.leadNIPR.text() #lead nipr
            lead_sipr = self.leadSIPR.text() #lead sipr
            lead_office = self.leadOffice.text() #lead office
            lead_mobile = self.leadMobile.text() #lead mobile
            office_symbol = self.office_comboBox.currentText()
            event_name = self.eventName.text() #event name
            start_date = self.startDate.text() #start date
            end_date = self.endDate.text() #end date
            if self.unclassButton.isChecked() == True:
                classification = 'UNCLASSIFIED' #classification - unclassified
                if self.cuiButton.isChecked() == True:
                    designation = 'CUI' #designation - cui
                elif self.fouoButton.isChecked() == True:
                    designation = 'FOUO' #designation - fouo
                else:
                    designation = 'NONE' #designation - none
            if self.secretButton.isChecked() == True:
                classification = 'SECRET' #classification - secret
                if self.noforncheckBox.isChecked() == True:
                    designation = 'NOFORN' #designation - noforn
                else:
                    designation = 'NONE' #designation - none
            if self.topsecretButton.isChecked() == True:
                classification = 'TOP SECRET' #classification - top secret
                if self.scicheckBox.isChecked() == True:
                    designation = 'SCI' #designation - sci
                else:
                    designation = 'NONE' #designation - none
            if self.draftcheckBox.isChecked() == True:
                draft = 'YES' #draft - yes
            else:
            	draft = 'NO' #draft - no
            if self.lightButton.isChecked() == True:
                mode = 'LIGHT' #mode - light
            if self.darkButton.isChecked() == True:
                mode = 'DARK' #mode - dark
            event = Event(lead_name, lead_title, lead_org, lead_nipr, lead_sipr, lead_office, lead_mobile, office_symbol, event_name, event_type, start_date, end_date, classification, designation, draft, mode)
            event.write_file()


        def go_page3():
            save_main()
            global event
            if event.classification == 'UNCLASSIFIED':
                if event.designation == 'CUI':
                    self.definedButton.setText('CUI')
                else:
                    self.definedButton.setText('UNCLASSIFIED//FOUO')
            if event.classification == 'SECRET':
                if event.designation == 'NOFORN':
                    self.definedButton.setText('SECRET//NOFORN')
                else:
                    self.definedButton.setText('SECRET')
            if event.classification == 'TOP SECRET':
                if event.designation == 'SCI':
                    self.definedButton.setText('TOP SECRET//SCI')
                else:
                    self.definedButton.setText('TOP SECRET')
            #SUGGEST DELIVERABLES BASED ON EVENT TYPE:
            sut_file = cwd + '/sut_info.xml'
            if os.path.exists(sut_file) == False:
                t_event = event.event_name.lower()
                if 'cvi' in t_event:
                    self.ERBcheckBox.setChecked(True)
                    self.ARMcheckBox.setChecked(True)
                    self.reportcheckBox.setChecked(False)
                elif 'vof' in t_event:
                    self.ERBcheckBox.setChecked(True)
                    self.ARMcheckBox.setChecked(False)
                    self.reportcheckBox.setChecked(True)
                    self.memoButton.setChecked(True)
                    self.inputDays.setText("30")
                else:
                    self.ERBcheckBox.setChecked(True)
                    self.ARMcheckBox.setChecked(True)
                    self.reportcheckBox.setChecked(True)
                    self.reportButton.setChecked(True)
                    self.inputDays.setText("90")
                deliverables_changed()
            self.stackedWidget0.setCurrentIndex(2)


        def go_page2():
            save_main()
            if self.cvpaButton.isChecked() == True:
                if '.zip' in data_source:
                    aux_data_folder = data_folder.replace(' ', '\ ')
                    msg = 'python3 dradis_parser.py ' + aux_data_folder
                    msg = msg.replace('(', '\(')
                    msg = msg.replace(')', '\)')
                    os.system(msg)
                if 'fric_export_' in data_source:
                    msg = 'python3 fric_parser.py ' + data_folder
                    os.system(msg)
                if 'Create your own ERB' in data_source:
                    msg = 'python3 empty_erb.py'
                    os.system(msg)
                current_selection = self.listFolders.currentItem().text()
                if folder_list != current_selection:
                    get_findings()
                self.stackedWidget0.setCurrentIndex(1)
                sutData()
            if self.pmrButton.isChecked() == True:
                create_pptx()


        def go_page1(wp):
            if wp == 1:
                global folder_list
                folder_list = self.listFolders.currentItem().text()
            self.stackedWidget0.setCurrentIndex(0)


        self.leadNIPR.textChanged.connect(nipr_changed)
        self.cvpaButton.toggled.connect(event_type)
        self.pmrButton.toggled.connect(event_type)
        self.systemButton.clicked.connect(go_page3)
        self.unclassButton.toggled.connect(class_changed)
        self.secretButton.toggled.connect(class_changed)
        self.topsecretButton.toggled.connect(class_changed)
        self.listFolders.itemSelectionChanged.connect(on_selection_changed)
        self.listFindings.itemSelectionChanged.connect(set_fields)
        self.listScreenshots.itemSelectionChanged.connect(preview_image)
        self.nextButton.clicked.connect(go_page2)
        self.gobackButton.clicked.connect(go_page1, 1)
        self.findingAddButton.clicked.connect(addFinding)
        self.findingUpButton.clicked.connect(findingsUp)
        self.findingDownButton.clicked.connect(findingsDown)
        self.findingDeleteButton.clicked.connect(delFinding)
        self.updateDescButton.clicked.connect(update_finding)
        self.screenshotAddButton.clicked.connect(addScreenshot)
        self.screenshotUpButton.clicked.connect(screenshotsUp)
        self.screenshotDownButton.clicked.connect(screenshotsDown)
        self.screenshotDeleteButton.clicked.connect(delScreenshot)
        self.pptxButton.clicked.connect(create_pptx)
        self.quitButton.clicked.connect(QApplication.instance().quit)
        self.scopeAddButton.clicked.connect(addScope)
        self.scopeDeleteButton.clicked.connect(delScope)
        self.reportcheckBox.toggled.connect(deliverables_changed)
        self.figureBrowseButton.clicked.connect(addDiagram)
        self.diagramDeleteButton.clicked.connect(delDiagram)
        self.diagramUpButton.clicked.connect(diagramsUp)
        self.diagramDownButton.clicked.connect(diagramsDown)
        self.listSUT_diagrams.itemSelectionChanged.connect(preview_SUT)
        self.strengthsAddButton.clicked.connect(addStrengths)
        self.strengthsDeleteButton.clicked.connect(delStrengths)
        self.strengthsUpButton.clicked.connect(strengthsUp)
        self.strengthsDownButton.clicked.connect(strengthsDown)
        self.weaknessesAddButton.clicked.connect(addWeaknesses)
        self.weaknessesDeleteButton.clicked.connect(delWeaknesses)
        self.weaknessesUpButton.clicked.connect(weaknessesUp)
        self.weaknessesDownButton.clicked.connect(weaknessesDown)
        self.mitigationsAddButton.clicked.connect(addMitigations)
        self.mitigationsDeleteButton.clicked.connect(delMitigations)
        self.mitigationsUpButton.clicked.connect(mitigationsUp)
        self.mitigationsDownButton.clicked.connect(mitigationsDown)
        self.saveButton.clicked.connect(sutData)
        self.saveButton.clicked.connect(go_page1, 2)


    def retranslateUi(self, MainWindow):
        _translate = QtCore.QCoreApplication.translate
        MainWindow.setWindowTitle(_translate("MainWindow", "Emerging Results Brief (ERB) Generator"))
        self.startDate.setDisplayFormat(_translate("MainWindow", "MM/dd/yyyy"))
        self.leadName_label.setText(_translate("MainWindow", "Name:"))
        self.leadName.setPlaceholderText(_translate("MainWindow", "Your name here!"))
        self.leadTitle.setText(_translate("MainWindow", "Computer Scientist"))
        self.eventName.setPlaceholderText(_translate("MainWindow", "Title of the Event"))
        self.leadTitle_label.setText(_translate("MainWindow", "Rank/Title:"))
        self.leadOrg_label.setText(_translate("MainWindow", "Organization:"))
        self.endDate.setDisplayFormat(_translate("MainWindow", "MM/dd/yyyy"))
        self.startDate_label.setText(_translate("MainWindow", "Start Date:"))
        self.endDate_label.setText(_translate("MainWindow", "End Date:"))
        self.dradisfricData_label.setText(_translate("MainWindow", "DRADIS/FRIC DATA"))
        self.eventName_label.setText(_translate("MainWindow", "Name:"))
        self.leadOrg.setText(_translate("MainWindow", "DEVCOM Analysis Center"))
        self.event_label.setText(_translate("MainWindow", "EVENT"))
        self.teamLead_label.setText(_translate("MainWindow", "TEAM LEAD"))
        self.availFilesFolders_label.setText(_translate("MainWindow", "Available Files / Folders:"))
        self.nextButton.setText(_translate("MainWindow", "Findings"))
        self.unclassButton.setText(_translate("MainWindow", "UNCLASSIFIED"))
        self.secretButton.setText(_translate("MainWindow", "SECRET"))
        self.topsecretButton.setText(_translate("MainWindow", "TOP SECRET"))
        self.noforncheckBox.setText(_translate("MainWindow", "NOFORN"))
        self.scicheckBox.setText(_translate("MainWindow", "SCI"))
        self.fouoButton.setText(_translate("MainWindow", "FOUO (obsolete)"))
        self.cuiButton.setText(_translate("MainWindow", "CUI"))
        self.mode_label.setText(_translate("MainWindow", "Slides\n"
"Background"))
        self.lightButton.setText(_translate("MainWindow", "LIGHT"))
        self.darkButton.setText(_translate("MainWindow", "DARK"))
        self.classification_label.setText(_translate("MainWindow", "Classification"))
        self.draftcheckBox.setText(_translate("MainWindow", "DRAFT//PRE-DECISIONAL"))
        self.office_label.setText(_translate("MainWindow", "Office Symbol:"))
        self.pmrButton.setText(_translate("MainWindow", "PMR"))
        self.cvpaButton.setText(_translate("MainWindow", "CVPA / CVI"))
        self.eventType_label.setText(_translate("MainWindow", "Type:"))
        self.leadEmail_label.setText(_translate("MainWindow", "Emails:"))
        self.leadNIPR.setText(_translate("MainWindow", "xxxxxxxx@army.mil"))
        self.leadNIPR_label.setText(_translate("MainWindow", "NIPR:"))
        self.leadSIPR.setText(_translate("MainWindow", "xxxxxxxx@mail.smil.mil"))
        self.leadSIPR_label.setText(_translate("MainWindow", "SIPR:"))
        self.leadPhone_label.setText(_translate("MainWindow", "Phones:"))
        self.leadOffice.setText(_translate("MainWindow", "(575) 678-xxxx"))
        self.leadMobile.setText(_translate("MainWindow", "(575) xxx-xxxx"))
        self.leadOffice_label.setText(_translate("MainWindow", "Office:"))
        self.leadMobile_label.setText(_translate("MainWindow", "Mobile:"))
        self.systemButton.setText(_translate("MainWindow", "System\n"
"Info"))
        self.quitButton.setText(_translate("MainWindow", "Quit"))
        self.findings_label.setText(_translate("MainWindow", "Findings:"))
        self.updateDescButton.setText(_translate("MainWindow", "Update Finding"))
        self.screenshotPreview_label.setText(_translate("MainWindow", "Screenshot Preview:"))
        self.imagePreview_label.setText(_translate("MainWindow", "NO PREVIEW"))
        self.findingDeleteButton.setText(_translate("MainWindow", "DEL"))
        self.findingUpButton.setText(_translate("MainWindow", "UP"))
        self.findingDownButton.setText(_translate("MainWindow", "DOWN"))
        self.screenshots_label.setText(_translate("MainWindow", "Screenshots:"))
        self.pptxButton.setText(_translate("MainWindow", "PPTX"))
        self.issues_label.setText(_translate("MainWindow", "Issues:"))
        self.gobackButton.setText(_translate("MainWindow", "Go Back"))
        self.screenshotDeleteButton.setText(_translate("MainWindow", "DEL"))
        self.screenshotDownButton.setText(_translate("MainWindow", "DOWN"))
        self.screenshotUpButton.setText(_translate("MainWindow", "UP"))
        self.findingName_label.setText(_translate("MainWindow", "Finding Name:"))
        self.findingHosts_label.setText(_translate("MainWindow", "Affected Hosts:"))
        self.posture_label.setText(_translate("MainWindow", "Posture:"))
        self.insiderButton.setText(_translate("MainWindow", "Insider"))
        self.nearsiderButton.setText(_translate("MainWindow", "Nearsider"))
        self.outsiderButton.setText(_translate("MainWindow", "Outsider"))
        self.mitigation_label.setText(_translate("MainWindow", "Mitigation:"))
        self.mitigationcheckBox.setText(_translate("MainWindow", "Include Mitigation"))
        self.findingAddButton.setText(_translate("MainWindow", "ADD"))
        self.screenshotAddButton.setText(_translate("MainWindow", "ADD"))
        self.strengthsDeleteButton.setText(_translate("MainWindow", "DEL"))
        self.strengths_label.setText(_translate("MainWindow", "Strengths:"))
        self.strengthsAddButton.setText(_translate("MainWindow", "ADD"))
        self.sutPreview_label.setText(_translate("MainWindow", "SUT Preview:"))
        self.sutPreview_image.setText(_translate("MainWindow", "NO PREVIEW"))
        self.SUT_label.setText(_translate("MainWindow", "SYSTEM UNDER TEST (SUT)"))
        self.weaknesses_label.setText(_translate("MainWindow", "Weaknesses:"))
        self.weaknessesAddButton.setText(_translate("MainWindow", "ADD"))
        self.weaknessesDeleteButton.setText(_translate("MainWindow", "DEL"))
        self.mitigationsDeleteButton.setText(_translate("MainWindow", "DEL"))
        self.mitigationsAddButton.setText(_translate("MainWindow", "ADD"))
        self.mitigations_label.setText(_translate("MainWindow", "Overall Mitigations:"))
        self.inputStrength.setPlaceholderText(_translate("MainWindow", "Type Strength to add"))
        self.inputStrength_label.setText(_translate("MainWindow", "Input:"))
        self.inputWeakness_label.setText(_translate("MainWindow", "Input:"))
        self.inputWeakness.setPlaceholderText(_translate("MainWindow", "Type Weakness to add"))
        self.inputMitigation_label.setText(_translate("MainWindow", "Input:"))
        self.inputMitigation.setPlaceholderText(_translate("MainWindow", "Type Mitigation to add"))
        self.saveButton.setText(_translate("MainWindow", "SAVE"))
        self.diagramSUT_label.setText(_translate("MainWindow", "Diagram/Figure :"))
        self.figureBrowseButton.setText(_translate("MainWindow", "BROWSE"))
        self.inputScope_label.setText(_translate("MainWindow", "Input:"))
        self.scope_label.setText(_translate("MainWindow", "Scope:"))
        self.scopeAddButton.setText(_translate("MainWindow", "ADD"))
        self.inputScope.setPlaceholderText(_translate("MainWindow", "Type Scope to add"))
        self.scopeDeleteButton.setText(_translate("MainWindow", "DEL"))
        self.subbulletButton.setText(_translate("MainWindow", "Sub-Bullet (indent)"))
        self.bulletButton.setText(_translate("MainWindow", "Bullet"))
        self.reportButton.setText(_translate("MainWindow", "Report"))
        self.memoButton.setText(_translate("MainWindow", "Memorandum"))
        self.inputDays.setText(_translate("MainWindow", "90"))
        self.inputDays.setPlaceholderText(_translate("MainWindow", "Add # of business days"))
        self.inputDays_label.setText(_translate("MainWindow", "business days"))
        self.deliverables_label.setText(_translate("MainWindow", "Deliverables:"))
        self.ARMcheckBox.setText(_translate("MainWindow", "ARM (Risk Matrix)"))
        self.ERBcheckBox.setText(_translate("MainWindow", "ERB"))
        self.markings_label.setText(_translate("MainWindow", "Classification Markings:"))
        self.definedButton.setText(_translate("MainWindow", "EVENT DEFINED"))
        self.unclassifiedButton.setText(_translate("MainWindow", "UNCLASSIFIED"))
        self.diagramDeleteButton.setText(_translate("MainWindow", "DEL"))
        self.listSUT_diagrams_label.setText(_translate("MainWindow", "List of Diagrams/Figures:"))
        self.diagramUpButton.setText(_translate("MainWindow", "UP"))
        self.diagramDownButton.setText(_translate("MainWindow", "DOWN"))
        self.strengthsUpButton.setText(_translate("MainWindow", "UP"))
        self.strengthsDownButton.setText(_translate("MainWindow", "DOWN"))
        self.weaknessesDownButton.setText(_translate("MainWindow", "DOWN"))
        self.weaknessesUpButton.setText(_translate("MainWindow", "UP"))
        self.mitigationsDownButton.setText(_translate("MainWindow", "DOWN"))
        self.mitigationsUpButton.setText(_translate("MainWindow", "UP"))
        self.reportcheckBox.setText(_translate("MainWindow", "Report/Memo"))


if __name__ == "__main__":
    import sys
    def exitHandler():
        try:
            #REMOVE INACTIVE FINDINGS IN XML
            tree = ET.parse(current_erb)
            root = tree.getroot()
            find_to_del = []
            for finding in root.findall('finding'):
                uid = finding.get('uid')
                active = int(finding.find('active').text)
                if active == 0:
                    root.remove(finding)
                    folder = finding.find('folder')
                    find_to_del.append(folder.text)
            tree.write(current_erb)
            #UPDATE AND SORT FINDINGS BY RANK
            tree = ET.parse(current_erb)
            root = tree.getroot()
            findings_dict = {}
            for finding in root.findall('finding'):
                uid = finding.get('uid')
                rank = int(finding.find('rank').text)
                findings_dict[uid] = rank
            sorted_dict = sorted(findings_dict.items(), key=lambda kv:kv[1])
            new_dict = {}
            for i in range(len(sorted_dict)):
                new_dict[sorted_dict[i][0]] = i
            for finding in root.findall('finding'):
                uid = finding.get('uid')
                rank = finding.find('rank')
                rank.text = str(new_dict[uid])
            tree.write(current_erb)
            #UPDATE UIDs
            tree = ET.parse(current_erb)
            root = tree.getroot()
            sorted_dict = sorted(new_dict.items(), key=lambda kv:int(kv[0]))
            new_dict = {}
            for i in range(len(sorted_dict)):
                new_dict[sorted_dict[i][0]] = sorted_dict[i][1]
            for finding in root.findall('finding'):
                uid = finding.get('uid')
                index = list(new_dict.keys()).index(uid)
                finding.set('uid', str(index))
            tree.write(current_erb)
            #DELETE UNUSED SCREENSHOTS
            for i in find_to_del:
                if os.path.exists(i):
                    os.system('rm -r ' + i)
        except:
            pass
    app = QtWidgets.QApplication(sys.argv)
    app.aboutToQuit.connect(exitHandler)
    MainWindow = QtWidgets.QMainWindow()
    ui = Ui_MainWindow()
    ui.setupUi(MainWindow)
    MainWindow.show()
    sys.exit(app.exec_())
